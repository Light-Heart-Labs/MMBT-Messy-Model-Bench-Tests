import csv
import textwrap
from pathlib import Path

from PIL import Image, ImageDraw, ImageFont
from pptx import Presentation
from pptx.util import Inches

ROOT = Path(__file__).resolve().parents[2]
CHARTS = ROOT / "assets/charts"
DIAGRAMS = ROOT / "assets/diagrams"
TABLES = ROOT / "assets/tables"
PREVIEWS = ROOT / "deck/previews"
DECK = ROOT / "deck/yeti_board_advisors_presentation.pptx"
PDF = ROOT / "deck/yeti_board_advisors_presentation.pdf"
INVENTORY = ROOT / "deck/source/slide_inventory.csv"

W, H = 1600, 900
PAPER = "#F7F4EC"
INK = "#17212B"
MUTED = "#607080"
GRID = "#D9E1E8"
BLUE = "#2A6FDB"
GREEN = "#27AE60"
RED = "#D64A4A"
AMBER = "#D99A2B"
TEAL = "#009CA6"
PURPLE = "#7A4CC2"
WHITE = "#FFFFFF"
TRACE_COMMIT = "8bb17db"


def font(size, bold=False):
    candidates = [
        "C:/Windows/Fonts/arialbd.ttf" if bold else "C:/Windows/Fonts/arial.ttf",
        "C:/Windows/Fonts/calibrib.ttf" if bold else "C:/Windows/Fonts/calibri.ttf",
        "/usr/share/fonts/truetype/dejavu/DejaVuSans-Bold.ttf" if bold else "/usr/share/fonts/truetype/dejavu/DejaVuSans.ttf",
        "/usr/share/fonts/dejavu/DejaVuSans-Bold.ttf" if bold else "/usr/share/fonts/dejavu/DejaVuSans.ttf",
    ]
    for path in candidates:
        if Path(path).exists():
            return ImageFont.truetype(path, size=size)
    return ImageFont.load_default()


F_HERO = font(86, True)
F_TITLE = font(50, True)
F_SUB = font(25)
F_BODY = font(24)
F_BODY_B = font(24, True)
F_SMALL = font(18)
F_SMALL_B = font(18, True)
F_TINY = font(14)
F_NUM = font(42, True)


def read_csv(name):
    with (TABLES / name).open(newline="", encoding="utf-8") as f:
        return list(csv.DictReader(f))


def fmt_money(value, decimals=0):
    return f"${float(value):,.{decimals}f}"


def fmt_pct(value, decimals=1):
    return f"{float(value) * 100:.{decimals}f}%"


def wrap_lines(text, width):
    return textwrap.wrap(text, width=width, break_long_words=False)


def draw_text_box(draw, xy, text, font_obj=F_BODY, fill=INK, width=54, line_gap=8):
    x, y = xy
    for line in wrap_lines(text, width):
        draw.text((x, y), line, font=font_obj, fill=fill)
        y += font_obj.size + line_gap
    return y


def rounded(draw, xy, fill=WHITE, outline=None, width=2, radius=18):
    draw.rounded_rectangle(xy, radius=radius, fill=fill, outline=outline, width=width)


def new_slide(title, subtitle="", accent=BLUE):
    im = Image.new("RGB", (W, H), PAPER)
    draw = ImageDraw.Draw(im)
    draw.rectangle((0, 0, W, 12), fill=accent)
    draw.text((64, 40), title, fill=INK, font=F_TITLE)
    if subtitle:
        draw.text((68, 102), subtitle, fill=MUTED, font=F_SUB)
    return im, draw


def footer(draw, slide_no, trace, extra=""):
    y = H - 46
    draw.line((64, y - 14, W - 64, y - 14), fill=GRID, width=2)
    text = f"Trace: audit/traces/slide-{slide_no:02d}.md | input {TRACE_COMMIT} | {trace}"
    if extra:
        text += f" | {extra}"
    draw.text((66, y), text, fill=MUTED, font=F_TINY)


def paste_fit(base, path, box, crop_top=0, bg=PAPER):
    img = Image.open(path).convert("RGB")
    if crop_top:
        img = img.crop((0, crop_top, img.width, img.height))
    x0, y0, x1, y1 = box
    bw, bh = x1 - x0, y1 - y0
    scale = min(bw / img.width, bh / img.height)
    nw, nh = int(img.width * scale), int(img.height * scale)
    resized = img.resize((nw, nh), Image.Resampling.LANCZOS)
    panel = Image.new("RGB", (bw, bh), bg)
    panel.paste(resized, ((bw - nw) // 2, (bh - nh) // 2))
    base.paste(panel, (x0, y0))
    return (x0, y0, x1, y1)


def card(draw, xy, title, body, color=BLUE, body_width=38, title_font=F_BODY_B):
    rounded(draw, xy, fill=WHITE, outline=color, width=3, radius=14)
    x0, y0, x1, y1 = xy
    draw.rectangle((x0, y0, x0 + 10, y1), fill=color)
    draw.text((x0 + 28, y0 + 22), title, fill=INK, font=title_font)
    draw_text_box(draw, (x0 + 28, y0 + 60), body, font_obj=F_SMALL, fill=MUTED, width=body_width, line_gap=6)


def bullet_list(draw, x, y, bullets, color=BLUE, width=46, gap=20):
    for title, body, trace in bullets:
        draw.ellipse((x, y + 7, x + 14, y + 21), fill=color)
        draw.text((x + 28, y), title, fill=INK, font=F_BODY_B)
        y = draw_text_box(draw, (x + 28, y + 34), body, font_obj=F_SMALL, fill=MUTED, width=width, line_gap=6)
        if trace:
            draw.text((x + 28, y + 4), trace, fill=color, font=F_TINY)
            y += 26
        y += gap
    return y


def render_chart_slide(slide_no, title, subtitle, image_path, trace, accent=BLUE, crop_top=128, callouts=None):
    im, draw = new_slide(title, subtitle, accent)
    paste_fit(im, image_path, (66, 150, 1070, 800), crop_top=crop_top)
    rounded(draw, (1100, 150, 1538, 800), fill=WHITE, outline=GRID, width=2, radius=18)
    if callouts:
        bullet_list(draw, 1130, 180, callouts, color=accent, width=28, gap=18)
    footer(draw, slide_no, trace)
    return im


def slide_01():
    rows = {r["metric"]: r for r in read_csv("recommendation_snapshot.csv")}
    im, draw = new_slide("Recommendation: HOLD", "Small upside means the point is restraint, not indecision.", BLUE)
    draw.text((76, 172), "YETI", fill=INK, font=font(42, True))
    draw.text((76, 222), "Hold the stock; do not add risk here.", fill=MUTED, font=F_SUB)
    draw.text((76, 305), "HOLD", fill=BLUE, font=F_HERO)
    metrics = [
        ("Current", fmt_money(rows["Current price"]["value"], 2), "T001", BLUE),
        ("Target", fmt_money(rows["12-month target"]["value"], 0), "T003", AMBER),
        ("Upside", fmt_pct(rows["Upside"]["value"], 1), "T004", GREEN),
        ("Market cap", "$3.0B", "T002", PURPLE),
    ]
    x = 520
    for label, value, trace, color in metrics:
        rounded(draw, (x, 198, x + 230, 340), fill=WHITE, outline=color, width=3)
        draw.text((x + 24, 226), value, fill=color, font=F_NUM)
        draw.text((x + 24, 282), label, fill=INK, font=F_SMALL_B)
        draw.text((x + 24, 308), trace, fill=MUTED, font=F_TINY)
        x += 250
    paste_fit(im, CHARTS / "scenario_distribution.png", (76, 405, 1518, 790), crop_top=126)
    draw.text((80, 808), "Confidence: moderate. Historical data are strong; valuation and tariff assumptions are judgment-heavy.", fill=INK, font=F_SMALL_B)
    footer(draw, 1, "T001-T004, T023-T026")
    return im


def slide_02():
    im, draw = new_slide("The Thesis In The Agent's Words", "Quality is real. The price already knows enough of the good news.", TEAL)
    thesis = "YETI is a durable premium brand with international and DTC runway, but the public evidence does not support paying for a clean upside inflection while tariff and US-category risk remain live."
    rounded(draw, (80, 170, 1518, 360), fill=WHITE, outline=TEAL, width=3)
    draw_text_box(draw, (118, 205), thesis, font_obj=font(33, True), fill=INK, width=78, line_gap=10)
    bullets = [
        ("Why not sell", "FY2025 free cash flow and gross margin remain healthy enough to avoid a bearish call.", "T007, T008, T010"),
        ("Why not buy", "The blended target is only modestly above the current price and tariffs remain the largest swing factor.", "T003, T004, T042"),
        ("What matters next", "International mix, DTC strength, tariff mitigation, and US wholesale sell-through decide whether HOLD becomes directional.", "T011-T015, Q1-Q3"),
    ]
    bullet_list(draw, 120, 430, bullets, color=TEAL, width=82, gap=24)
    footer(draw, 2, "T003-T015, T042, Q1-Q3")
    return im


def slide_03():
    im, draw = new_slide("What Would Change The Recommendation", "The recommendation is falsifiable, not a one-way story.", AMBER)
    paste_fit(im, DIAGRAMS / "decision_tree.png", (76, 150, 1060, 795), crop_top=128)
    callouts = [
        ("Upgrade path", "Tariff relief arrives faster than modeled, international growth stays above the US mix, and sell-through closes the wholesale gap.", "T014, T042, Q2"),
        ("Downgrade path", "Tariff costs persist, category demand weakens, or wholesale caution becomes real sell-through deterioration.", "T012, T042, T043"),
        ("Stay HOLD", "Execution is credible but upside fails the 20% risk-adjusted hurdle in the valuation ADR.", "/input/repo/decisions/0005"),
    ]
    rounded(draw, (1092, 150, 1536, 795), fill=WHITE, outline=GRID, width=2)
    bullet_list(draw, 1124, 185, callouts, color=AMBER, width=28, gap=22)
    footer(draw, 3, "T012, T014, T042-T043, Q2")
    return im


def slide_04():
    return render_chart_slide(
        4,
        "Evidence Stack",
        "The agent built the recommendation from filings, transcripts, other primary sources, trace rows, ADRs, and commits.",
        CHARTS / "evidence_counts.png",
        "evidence_counts.csv and input repo inventory",
        PURPLE,
        crop_top=128,
        callouts=[
            ("Source breadth", "The input repo preserved filings, transcripts, market pages, release pages, extraction scripts, and SHA logs.", "audit/numbers.md"),
            ("Trace density", "Every investment number in the memo maps to a trace row before it appears in this deck.", "T000-T043"),
            ("Process evidence", "ADRs and commit history are treated as reasoning artifacts, not housekeeping.", "/input/repo/decisions"),
        ],
    )


def slide_05():
    return render_chart_slide(
        5,
        "Financial Trajectory",
        "Guide-derived recovery is visible, but it is not a heroic margin rebound story.",
        CHARTS / "financial_trajectory.png",
        "T005-T010, T015-T017",
        BLUE,
        crop_top=128,
        callouts=[
            ("FY2025 base", "Revenue, operating margin, gross margin, and FCF all trace to the FY2025 10-K extraction.", "T005-T010"),
            ("FY2026 step", "The model uses the management guide-derived revenue growth assumption.", "T015-T017"),
            ("Inflection", "The trajectory supports quality, but the margin recovery is gradual enough to cap conviction.", "T009, T015"),
        ],
    )


def slide_06():
    im = render_chart_slide(
        6,
        "Business Mix",
        "DTC and international are real, but not a free pass through US/category risk.",
        CHARTS / "business_mix.png",
        "T011-T014, Q1",
        GREEN,
        crop_top=128,
        callouts=[
            ("DTC scale", "DTC sales are a majority of FY2025 sales, which supports data access and margin potential.", "T011"),
            ("Wholesale still matters", "Wholesale remains material, so channel caution cannot be ignored.", "T012, Q2"),
            ("International runway", "Management frames international as a long-term growth driver; the deck quotes it with context.", "Q1 lines 196-210"),
        ],
    )
    draw = ImageDraw.Draw(im)
    rounded(draw, (1120, 662, 1518, 783), fill="#F0FAF6", outline=GREEN, width=2)
    draw_text_box(draw, (1144, 682), '"international addressable market exceeds the US"', font_obj=F_SMALL_B, fill=INK, width=35)
    draw.text((1144, 746), "Q1, line 199; context in audit/quotes.md", fill=MUTED, font=F_TINY)
    return im


def slide_07():
    return render_chart_slide(
        7,
        "Valuation Bridge",
        "The price target comes from a blended valuation, with DCF anchoring the call.",
        CHARTS / "valuation_bridge.png",
        "T001, T003, T020-T022",
        AMBER,
        crop_top=128,
        callouts=[
            ("DCF anchor", "The DCF value is below the blended target, which tempers multiple-driven optimism.", "T020"),
            ("Multiples cross-check", "EV/EBITDA and P/E support a higher cross-check but do not justify a BUY alone.", "T021-T022"),
            ("Recommendation rule", "The valuation ADR required meaningful risk-adjusted upside before moving to BUY.", "/input/repo/decisions/0005"),
        ],
    )


def slide_08():
    return render_chart_slide(
        8,
        "Scenario Distribution",
        "Bear/base/bull are weighted outcomes, not three decorative bullets.",
        CHARTS / "scenario_distribution.png",
        "T001, T003, T023-T026",
        PURPLE,
        crop_top=128,
        callouts=[
            ("Bear case", "Tariff and demand pressure move value below current price.", "T023"),
            ("Base case", "Execution recovery creates upside, but not enough for a BUY after blending.", "T024"),
            ("Bull case", "A stronger recovery exists, but the probability weight keeps the call restrained.", "T025-T026"),
        ],
    )


def slide_09():
    return render_chart_slide(
        9,
        "Competitive Position",
        "YETI screens premium, but the peer context does not scream cheap.",
        CHARTS / "peer_context.png",
        "T031-T034",
        TEAL,
        crop_top=128,
        callouts=[
            ("Premium gross margin", "YETI's gross margin supports brand quality relative to the public peer set.", "T010, T031"),
            ("Operating context", "The operating-margin gap is not enough to call the stock obviously mispriced.", "T031-T034"),
            ("Comp humility", "No public peer is perfect; the ADR makes the imperfection explicit.", "/input/repo/decisions/0003"),
        ],
    )


def slide_10():
    im, draw = new_slide("Why These Comps", "The peer set is a context lens, not a mechanical answer.", TEAL)
    rows = read_csv("comp_rationale.csv")
    x_positions = [74, 456, 838, 1220]
    colors = [GREEN, AMBER, RED, MUTED]
    for row, x, color in zip(rows, x_positions, colors):
        rounded(draw, (x, 185, x + 320, 660), fill=WHITE, outline=color, width=3)
        draw.text((x + 28, 220), row["ticker"], fill=color, font=font(52, True))
        draw.text((x + 28, 286), row["role"], fill=INK, font=F_SMALL_B)
        draw_text_box(draw, (x + 28, 340), row["rationale"], font_obj=F_SMALL, fill=MUTED, width=28, line_gap=8)
    draw.text((82, 720), "Decision logic: include imperfect but useful public analogs; exclude Solo Brands from core triangulation after it distorted the comparison.", fill=INK, font=F_BODY_B)
    footer(draw, 10, "/input/repo/decisions/0003-peer-set.md")
    return im


def slide_11():
    return render_chart_slide(
        11,
        "Mispricing Test",
        "The agent could not support a durable sell-side blind spot.",
        CHARTS / "mispricing_test.png",
        "T027-T030, T035-T036",
        RED,
        crop_top=128,
        callouts=[
            ("No forced miss", "Consensus targets and estimates were not obviously too low relative to the agent's model.", "T027-T030"),
            ("Guide tension", "Consensus EPS sat above the high end of management's guide-derived adjusted EPS range.", "T030, T036"),
            ("Conclusion", "The stronger claim is efficient pricing: market skepticism offsets visible international growth.", "/input/repo/analysis/sell_side_gap.md"),
        ],
    )


def slide_12():
    return render_chart_slide(
        12,
        "Risk Assessment",
        "Risks are ranked by investment impact and tied to what would have to be true.",
        CHARTS / "risk_register.png",
        "T012, T039, T042-T043",
        RED,
        crop_top=128,
        callouts=[
            ("Top risk", "Tariffs are the largest near-term swing factor because management quantified gross-margin pressure.", "T042-T043"),
            ("Demand/channel", "Wholesale caution matters because the channel remains material to sales.", "T012, Q2"),
            ("Capital allocation", "Repurchases help, but they cannot carry the thesis alone.", "T039"),
        ],
    )


def slide_13():
    im, draw = new_slide("Reasoning Trail", "From raw evidence to extraction, model, ADRs, dead ends, and conclusion.", PURPLE)
    paste_fit(im, DIAGRAMS / "reasoning_graph.png", (60, 148, 1540, 790), crop_top=128)
    footer(draw, 13, "/input/repo/raw, /input/repo/extracted, /input/repo/model, /input/repo/decisions")
    return im


def slide_14():
    return render_chart_slide(
        14,
        "Commit History As Reasoning Evidence",
        "The agent's history shows staged inquiry, not one-shot slide polish.",
        CHARTS / "commit_history.png",
        "input git log at 8bb17db",
        BLUE,
        crop_top=128,
        callouts=[
            ("Early structure", "The memo repo committed source registry, raw inputs, extraction, model, and memo in stages.", "input git log"),
            ("Why it matters", "Commit messages make reasoning transitions auditable, especially when a thesis changes.", "README + git log"),
            ("Board read", "The process artifact is evidence about the agent system, not merely about YETI.", "narrative/audience-analysis.md"),
        ],
    )


def slide_15():
    im, draw = new_slide("Dead Ends The Agent Rejected", "The rejected paths matter because they reveal what the system refused to overclaim.", AMBER)
    rows = read_csv("dead_ends.csv")
    y = 178
    colors = [BLUE, RED, PURPLE]
    for row, color in zip(rows, colors):
        rounded(draw, (90, y, 1510, y + 155), fill=WHITE, outline=color, width=3)
        draw.text((124, y + 28), row["hypothesis"], fill=INK, font=F_BODY_B)
        draw_text_box(draw, (124, y + 66), row["finding"], font_obj=F_BODY, fill=MUTED, width=88, line_gap=8)
        draw.text((124, y + 124), row["trace"], fill=color, font=F_TINY)
        y += 190
    draw.text((96, 762), "Board interpretation: this is a HOLD because the agent found support for quality, but not enough support for a mispricing claim.", fill=INK, font=F_BODY_B)
    footer(draw, 15, "/input/repo/research/dead-ends.md")
    return im


def slide_16():
    im, draw = new_slide("Traceability System", "The audit mechanism is a first-class part of the deliverable.", GREEN)
    paste_fit(im, DIAGRAMS / "audit_workflow.png", (70, 150, 1530, 788), crop_top=128)
    footer(draw, 16, "audit_steps.csv, audit/numbers.md, audit/traces/")
    return im


def slide_17():
    im, draw = new_slide("How To Verify This Deck", "Self-audit: pick a claim and reconstruct it from repo files.", GREEN)
    steps = read_csv("audit_steps.csv")
    x = 86
    y = 182
    for idx, row in enumerate(steps, start=1):
        color = [BLUE, TEAL, AMBER, PURPLE, GREEN][idx - 1]
        rounded(draw, (x, y, x + 282, y + 350), fill=WHITE, outline=color, width=3)
        draw.ellipse((x + 24, y + 24, x + 78, y + 78), fill=color)
        draw.text((x + 42, y + 35), str(idx), fill=WHITE, font=F_BODY_B)
        draw_text_box(draw, (x + 24, y + 104), row["action"], font_obj=F_BODY_B, fill=INK, width=21, line_gap=8)
        draw_text_box(draw, (x + 24, y + 204), row["artifact"], font_obj=F_SMALL, fill=MUTED, width=25, line_gap=6)
        x += 302
    rounded(draw, (88, 624, 1510, 778), fill="#F0FAF6", outline=GREEN, width=2)
    draw.text((120, 654), "Worked example:", fill=GREEN, font=F_BODY_B)
    draw_text_box(draw, (120, 694), "Slide 1 target $41 -> T003 -> audit/numbers.md -> /input/repo/analysis/memo_trace_table.csv:5 -> model/yeti_investment_model.xlsx Valuation!B19.", font_obj=F_BODY, fill=INK, width=94, line_gap=8)
    footer(draw, 17, "audit/traces/slide-01.md, audit/numbers.md, T003")
    return im


def slide_18():
    im, draw = new_slide("Reconciliation Spot Check", "Five deck numbers were traced end-to-end back into the input repo.", PURPLE)
    rows = read_csv("reconciliation.csv")
    headers = ["Number", "Trace", "Path"]
    xs = [90, 375, 520]
    y = 174
    draw.rectangle((76, y, 1524, y + 52), fill=PURPLE)
    for h, x in zip(headers, xs):
        draw.text((x, y + 14), h, fill=WHITE, font=F_SMALL_B)
    y += 64
    for idx, row in enumerate(rows):
        bg = WHITE if idx % 2 == 0 else "#EEF2F5"
        draw.rectangle((76, y - 8, 1524, y + 86), fill=bg)
        draw_text_box(draw, (90, y), row["number"], font_obj=F_SMALL_B, fill=INK, width=24, line_gap=5)
        draw.text((375, y), row["trace"], fill=PURPLE, font=F_SMALL_B)
        draw_text_box(draw, (520, y), row["source_chain"], font_obj=F_TINY, fill=MUTED, width=96, line_gap=4)
        y += 96
    draw.text((86, 760), "Full reconciliation narrative: audit/reconciliation.md", fill=INK, font=F_BODY_B)
    footer(draw, 18, "audit/reconciliation.md and reconciliation.csv")
    return im


def slide_19():
    return render_chart_slide(
        19,
        "Confidence Map",
        "The deck distinguishes sourced facts from judgment calls.",
        CHARTS / "confidence_map.png",
        "T005-T043, /input/repo/decisions",
        GREEN,
        crop_top=128,
        callouts=[
            ("Highest confidence", "Historical financials, market inputs, and guidance are directly source-backed.", "audit/numbers.md"),
            ("Medium confidence", "Scenario outcomes depend on modeled recovery and tariff assumptions.", "T023-T026, T042"),
            ("Lowest confidence", "A true sell-side edge would require paid models and channel checks.", "/input/repo/memo lines 53-55"),
        ],
    )


def slide_20():
    im, draw = new_slide("Limitations And Next Analyst Work", "End where the memo ends: confidence is strongest in facts, weaker in judgment.", RED)
    rows = read_csv("limitations_matrix.csv")
    x_positions = [76, 450, 824, 1198]
    colors = [GREEN, BLUE, AMBER, RED]
    for row, x, color in zip(rows, x_positions, colors):
        rounded(draw, (x, 175, x + 326, 710), fill=WHITE, outline=color, width=3)
        draw.text((x + 24, 212), row["category"], fill=color, font=F_BODY_B)
        draw.text((x + 24, 250), f"Confidence: {row['confidence']}", fill=INK, font=F_SMALL_B)
        draw_text_box(draw, (x + 24, 304), row["what_the_agent_knows"], font_obj=F_SMALL, fill=MUTED, width=29, line_gap=6)
        draw.line((x + 24, 445, x + 292, 445), fill=GRID, width=2)
        draw.text((x + 24, 468), "Human next step", fill=INK, font=F_SMALL_B)
        draw_text_box(draw, (x + 24, 500), row["what_a_human_would_add"], font_obj=F_SMALL, fill=MUTED, width=29, line_gap=6)
    draw.text((82, 764), "Final board takeaway: this is a capability demonstration of traceable restraint, not a claim that automation found a hidden trade.", fill=INK, font=F_BODY_B)
    footer(draw, 20, "/input/repo/memo/yeti_investment_memo.md lines 51-55")
    return im


SLIDES = [
    ("Recommendation: HOLD", slide_01),
    ("Thesis", slide_02),
    ("Recommendation Triggers", slide_03),
    ("Evidence Stack", slide_04),
    ("Financial Trajectory", slide_05),
    ("Business Mix", slide_06),
    ("Valuation Bridge", slide_07),
    ("Scenario Distribution", slide_08),
    ("Competitive Position", slide_09),
    ("Comp Rationale", slide_10),
    ("Mispricing Test", slide_11),
    ("Risk Assessment", slide_12),
    ("Reasoning Trail", slide_13),
    ("Commit History", slide_14),
    ("Dead Ends", slide_15),
    ("Traceability System", slide_16),
    ("How To Verify", slide_17),
    ("Reconciliation", slide_18),
    ("Confidence Map", slide_19),
    ("Limitations", slide_20),
]


def render_previews():
    PREVIEWS.mkdir(parents=True, exist_ok=True)
    paths = []
    inventory_rows = []
    for idx, (title, renderer) in enumerate(SLIDES, start=1):
        im = renderer()
        path = PREVIEWS / f"slide-{idx:02d}.png"
        im.save(path)
        paths.append(path)
        inventory_rows.append({"slide": idx, "title": title, "preview": str(path.relative_to(ROOT)), "trace": f"audit/traces/slide-{idx:02d}.md"})
    with INVENTORY.open("w", newline="", encoding="utf-8") as f:
        writer = csv.DictWriter(f, fieldnames=["slide", "title", "preview", "trace"])
        writer.writeheader()
        writer.writerows(inventory_rows)
    return paths


def build_pptx(preview_paths):
    prs = Presentation()
    prs.slide_width = Inches(13.333333)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]
    for path in preview_paths:
        slide = prs.slides.add_slide(blank)
        slide.shapes.add_picture(str(path), 0, 0, width=prs.slide_width, height=prs.slide_height)
    DECK.parent.mkdir(parents=True, exist_ok=True)
    prs.save(DECK)


def build_pdf(preview_paths):
    images = [Image.open(path).convert("RGB") for path in preview_paths]
    images[0].save(PDF, save_all=True, append_images=images[1:], resolution=144)


def build_contact_sheet(preview_paths):
    thumb_w, thumb_h = 320, 180
    cols = 4
    rows = (len(preview_paths) + cols - 1) // cols
    sheet = Image.new("RGB", (cols * thumb_w, rows * (thumb_h + 28)), PAPER)
    draw = ImageDraw.Draw(sheet)
    for idx, path in enumerate(preview_paths, start=1):
        img = Image.open(path).convert("RGB").resize((thumb_w, thumb_h), Image.Resampling.LANCZOS)
        col = (idx - 1) % cols
        row = (idx - 1) // cols
        x = col * thumb_w
        y = row * (thumb_h + 28)
        sheet.paste(img, (x, y + 28))
        draw.text((x + 10, y + 5), f"Slide {idx:02d}", fill=INK, font=F_TINY)
    sheet.save(PREVIEWS / "contact-sheet.png")


def main():
    preview_paths = render_previews()
    build_pptx(preview_paths)
    build_pdf(preview_paths)
    build_contact_sheet(preview_paths)
    print(f"Built {DECK}")
    print(f"Built {PDF}")
    print(f"Built {len(preview_paths)} preview images in {PREVIEWS}")


if __name__ == "__main__":
    main()
