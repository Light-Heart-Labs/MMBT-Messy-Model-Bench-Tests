#!/usr/bin/env python3
"""
Generate Board Presentation PowerPoint Deck
Generates: GTLB board presentation with all slides
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
import os

# Colors (matching ADR-005)
COLORS = {
    'primary': RGBColor(30, 136, 229),      # GitLab blue #1E88E5
    'secondary': RGBColor(67, 160, 71),     # Green #43A047
    'tertiary': RGBColor(229, 57, 53),      # Red #E53935
    'text_dark': RGBColor(33, 33, 33),      # Black #212121
    'text_light': RGBColor(255, 255, 255),  # White #FFFFFF
    'bg': RGBColor(250, 250, 250),          # Off-white #FAFAFA
    'grid': RGBColor(224, 224, 224),        # Light gray #E0E0E0
    'bear': RGBColor(211, 47, 47),          # Red #D32F2F
    'base': RGBColor(56, 142, 60),          # Green #388E3C
    'bull': RGBColor(25, 118, 210),         # Blue #1976D2
    'weighted': RGBColor(30, 136, 229),     # GitLab blue #1E88E5
}

# Create presentation
prs = Presentation()
prs.slide_width = Inches(10)
prs.slide_height = Inches(7.5)

def add_title_slide(prs, title, subtitle):
    """Add title slide"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
    
    # Background
    background = slide.background
    background.fill.solid()
    background.fill.fore_color.rgb = COLORS['bg']
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(1), Inches(2.5), Inches(8), Inches(1))
    title_frame = title_box.text_frame
    title_frame.text = title
    title_para = title_frame.paragraphs[0]
    title_para.font.size = Pt(44)
    title_para.font.bold = True
    title_para.font.color.rgb = COLORS['primary']
    title_para.alignment = PP_ALIGN.CENTER
    
    # Subtitle
    subtitle_box = slide.shapes.add_textbox(Inches(1), Inches(3.5), Inches(8), Inches(0.5))
    subtitle_frame = subtitle_box.text_frame
    subtitle_frame.text = subtitle
    subtitle_para = subtitle_frame.paragraphs[0]
    subtitle_para.font.size = Pt(24)
    subtitle_para.font.color.rgb = COLORS['text_dark']
    subtitle_para.alignment = PP_ALIGN.CENTER

def add_recommendation_slide(prs, recommendation, target, current, upside):
    """Add recommendation summary slide"""
    slide = prs.slides.add_slide(prs.slide_layouts[1])  # Title and Content
    
    # Title
    title_shape = slide.shapes.title
    title_shape.text = "Recommendation Summary"
    title_shape.text_frame.paragraphs[0].font.size = Pt(36)
    title_shape.text_frame.paragraphs[0].font.bold = True
    title_shape.text_frame.paragraphs[0].font.color.rgb = COLORS['primary']
    
    # Content
    body_shape = slide.placeholders[1]
    text_frame = body_shape.text_frame
    text_frame.clear()
    
    # Recommendation
    p = text_frame.paragraphs[0]
    p.text = f"**RECOMMENDATION:** {recommendation}"
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = COLORS['secondary']
    
    # Price target
    p = text_frame.add_paragraph()
    p.text = f"12-Month Price Target: ${target} (Upside: {upside})"
    p.font.size = Pt(28)
    p.font.color.rgb = COLORS['primary']
    
    # Current price
    p = text_frame.add_paragraph()
    p.text = f"Current Price: ${current}"
    p.font.size = Pt(24)
    p.font.color.rgb = COLORS['text_dark']
    
    # Probability-weighted target
    p = text_frame.add_paragraph()
    p.text = "Probability-Weighted Target: $45.98 (114% upside)"
    p.font.size = Pt(24)
    p.font.color.rgb = COLORS['weighted']

def add_thesis_slide(prs, thesis_points):
    """Add investment thesis slide"""
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    
    title_shape = slide.shapes.title
    title_shape.text = "Investment Thesis"
    title_shape.text_frame.paragraphs[0].font.size = Pt(36)
    title_shape.text_frame.paragraphs[0].font.bold = True
    title_shape.text_frame.paragraphs[0].font.color.rgb = COLORS['primary']
    
    body_shape = slide.placeholders[1]
    text_frame = body_shape.text_frame
    text_frame.clear()
    
    for i, point in enumerate(thesis_points):
        if i > 0:
            p = text_frame.add_paragraph()
        else:
            p = text_frame.paragraphs[0]
        
        p.text = point
        p.font.size = Pt(24)
        p.font.color.rgb = COLORS['text_dark']
        p.level = 0

def add_financial_trajectory_slide(prs):
    """Add financial trajectory slide"""
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    
    title_shape = slide.shapes.title
    title_shape.text = "Financial Trajectory: Historical + Projected"
    title_shape.text_frame.paragraphs[0].font.size = Pt(36)
    title_shape.text_frame.paragraphs[0].font.bold = True
    title_shape.text_frame.paragraphs[0].font.color.rgb = COLORS['primary']
    
    body_shape = slide.placeholders[1]
    text_frame = body_shape.text_frame
    text_frame.clear()
    
    # Key metrics
    metrics = [
        "Revenue: $424M (FY2023) → $955M (FY2026), 30% CAGR",
        "Operating Margins: -49.8% (FY2023) → -7.4% (FY2026), +13pp annually",
        "Free Cash Flow: $222M (FY2026), positive since FY2024",
        "ARR: $300M (FY2022) → $860M (FY2026), 30% CAGR",
        "NRR: 120% (FY2022) → 115% (FY2026), still strong",
    ]
    
    for metric in metrics:
        p = text_frame.add_paragraph()
        p.text = metric
        p.font.size = Pt(22)
        p.font.color.rgb = COLORS['text_dark']

def add_competitive_landscape_slide(prs):
    """Add competitive landscape slide"""
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    
    title_shape = slide.shapes.title
    title_shape.text = "Competitive Landscape"
    title_shape.text_frame.paragraphs[0].font.size = Pt(36)
    title_shape.text_frame.paragraphs[0].font.bold = True
    title_shape.text_frame.paragraphs[0].font.color.rgb = COLORS['primary']
    
    body_shape = slide.placeholders[1]
    text_frame = body_shape.text_frame
    text_frame.clear()
    
    # Key competitors
    competitors = [
        "GitLab (GTLB): 2.56x EV/Revenue, 23.2% growth, 87.4% gross margin",
        "Atlassian (TEAM): 3.22x EV/Revenue, 23.3% growth",
        "Datadog (DDOG): 12.39x EV/Revenue, 29.2% growth",
        "MongoDB (MDB): 7.33x EV/Revenue, 26.7% growth",
        "Peer Average: 6.71x EV/Revenue, 22.1% growth",
    ]
    
    for competitor in competitors:
        p = text_frame.add_paragraph()
        p.text = competitor
        p.font.size = Pt(22)
        p.font.color.rgb = COLORS['text_dark']

def add_valuation_slide(prs):
    """Add valuation approach slide"""
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    
    title_shape = slide.shapes.title
    title_shape.text = "Valuation Approach"
    title_shape.text_frame.paragraphs[0].font.size = Pt(36)
    title_shape.text_frame.paragraphs[0].font.bold = True
    title_shape.text_frame.paragraphs[0].font.color.rgb = COLORS['primary']
    
    body_shape = slide.placeholders[1]
    text_frame = body_shape.text_frame
    text_frame.clear()
    
    # Valuation methods
    methods = [
        "Primary: DCF with FCF projections",
        "  - WACC: 10.5%",
        "  - Terminal growth: 3.0%",
        "  - DCF base case: $38.52/share",
        "Secondary: EV/Revenue peer comparison",
        "  - Peer average: 6.71x",
        "  - Implied price: $53.25/share",
        "Selected target: $42.00 (midpoint, 3.5x FY2027E EV/Revenue)",
    ]
    
    for method in methods:
        p = text_frame.add_paragraph()
        p.text = method
        p.font.size = Pt(20)
        p.font.color.rgb = COLORS['text_dark']
        if method.startswith("  -"):
            p.level = 1

def add_risk_assessment_slide(prs):
    """Add risk assessment slide"""
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    
    title_shape = slide.shapes.title
    title_shape.text = "Risk Assessment"
    title_shape.text_frame.paragraphs[0].font.size = Pt(36)
    title_shape.text_frame.paragraphs[0].font.bold = True
    title_shape.text_frame.paragraphs[0].font.color.rgb = COLORS['primary']
    
    body_shape = slide.placeholders[1]
    text_frame = body_shape.text_frame
    text_frame.clear()
    
    # Risks
    risks = [
        "High Probability / High Impact:",
        "  • NRR decline (120% → 115%)",
        "  • Competition from GitHub Copilot",
        "Medium Probability / High Impact:",
        "  • Macro downturn (enterprise software spending cuts)",
        "  • AI monetization failure",
        "Low Probability / High Impact:",
        "  • Open-source cannibalization",
        "  • Key person risk (CEO Sid Sijbrandij)",
    ]
    
    for risk in risks:
        p = text_frame.add_paragraph()
        p.text = risk
        p.font.size = Pt(20)
        p.font.color.rgb = COLORS['text_dark']
        if risk.startswith("  •"):
            p.level = 1

def add_scenarios_slide(prs):
    """Add scenario analysis slide"""
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    
    title_shape = slide.shapes.title
    title_shape.text = "Bear/Base/Bull Scenarios"
    title_shape.text_frame.paragraphs[0].font.size = Pt(36)
    title_shape.text_frame.paragraphs[0].font.bold = True
    title_shape.text_frame.paragraphs[0].font.color.rgb = COLORS['primary']
    
    body_shape = slide.placeholders[1]
    text_frame = body_shape.text_frame
    text_frame.clear()
    
    # Scenarios
    scenarios = [
        "Bear (25% probability): $18.58/share (-14% downside)",
        "Base (50% probability): $38.52/share (79% upside)",
        "Bull (25% probability): $88.30/share (310% upside)",
        "",
        "Probability-Weighted Target: $45.98/share (114% upside)",
        "",
        "Current Price: $21.51",
    ]
    
    for scenario in scenarios:
        p = text_frame.add_paragraph()
        p.text = scenario
        p.font.size = Pt(22)
        p.font.color.rgb = COLORS['text_dark']
        if "Probability-Weighted" in scenario:
            p.font.bold = True
            p.font.color.rgb = COLORS['weighted']

def add_reasoning_trail_slide(prs):
    """Add reasoning trail slide"""
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    
    title_shape = slide.shapes.title
    title_shape.text = "Agent's Decision Process"
    title_shape.text_frame.paragraphs[0].font.size = Pt(36)
    title_shape.text_frame.paragraphs[0].font.bold = True
    title_shape.text_frame.paragraphs[0].font.color.rgb = COLORS['primary']
    
    body_shape = slide.placeholders[1]
    text_frame = body_shape.text_frame
    text_frame.clear()
    
    # Reasoning trail
    trail = [
        "INPUT DATA:",
        "  • SEC filings (yfinance)",
        "  • Earnings transcripts (Seeking Alpha)",
        "  • Financial data (yfinance)",
        "",
        "ANALYSIS PHASE:",
        "  • Financial modeling (three-statement model)",
        "  • Competitive analysis (8 comparable companies)",
        "  • Scenario analysis (bear/base/bull)",
        "",
        "DECISION PHASE:",
        "  • Investment thesis (5 key points)",
        "  • Valuation target ($42.00)",
        "  • Risk assessment (6 key risks)",
        "",
        "CONCLUSION:",
        "  • BUY recommendation",
        "  • 95% upside from current price",
    ]
    
    for step in trail:
        p = text_frame.add_paragraph()
        p.text = step
        p.font.size = Pt(18)
        p.font.color.rgb = COLORS['text_dark']
        if step.startswith("  •") or step.startswith("  -"):
            p.level = 1

def add_dead_ends_slide(prs):
    """Add dead ends slide"""
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    
    title_shape = slide.shapes.title
    title_shape.text = "Dead Ends: What Didn't Work"
    title_shape.text_frame.paragraphs[0].font.size = Pt(36)
    title_shape.text_frame.paragraphs[0].font.bold = True
    title_shape.text_frame.paragraphs[0].font.color.rgb = COLORS['primary']
    
    body_shape = slide.placeholders[1]
    text_frame = body_shape.text_frame
    text_frame.clear()
    
    # Dead ends
    dead_ends = [
        "1. SEC Filing Downloads",
        "   • Attempted: Direct URL access, session-based requests, curl",
        "   • Failed: 403 Forbidden errors",
        "   • Resolution: Used yfinance instead",
        "",
        "2. GitLab Investor Relations Website",
        "   • Attempted: Access investors.gitlab.com",
        "   • Failed: DNS resolution failed",
        "   • Resolution: Used Seeking Alpha transcripts",
        "",
        "3. PRNewswire Press Releases",
        "   • Attempted: Search PRNewswire",
        "   • Failed: Search URL returned 404",
        "   • Resolution: Used yfinance data",
        "",
        "4. Initial CIK Lookup",
        "   • Attempted: Used CIK 0001687226",
        "   • Failed: Mapped to wrong company",
        "   • Resolution: Found correct CIK 0001653482",
    ]
    
    for dead_end in dead_ends:
        p = text_frame.add_paragraph()
        p.text = dead_end
        p.font.size = Pt(18)
        p.font.color.rgb = COLORS['text_dark']
        if dead_end.startswith("   •") or dead_end.startswith("   -"):
            p.level = 1

def add_verification_slide(prs):
    """Add verification slide"""
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    
    title_shape = slide.shapes.title
    title_shape.text = "How to Audit This Deck"
    title_shape.text_frame.paragraphs[0].font.size = Pt(36)
    title_shape.text_frame.paragraphs[0].font.bold = True
    title_shape.text_frame.paragraphs[0].font.color.rgb = COLORS['primary']
    
    body_shape = slide.placeholders[1]
    text_frame = body_shape.text_frame
    text_frame.clear()
    
    # Verification steps
    verification = [
        "Every claim is traceable to source:",
        "",
        "1. TRACE FILES",
        "   Location: /audit/traces/",
        "   Purpose: Point to source files in input repo",
        "",
        "2. NUMBER TRACKING",
        "   Location: /audit/numbers.md",
        "   Purpose: Every number with source path",
        "",
        "3. QUOTE TRACKING",
        "   Location: /audit/quotes.md",
        "   Purpose: Every quote with file + line number",
        "",
        "4. RECONCILIATION",
        "   Location: /audit/reconciliation.md",
        "   Purpose: Spot-check 5 random numbers",
        "",
        "To verify any claim: Follow the trace file → Check source → Verify number",
    ]
    
    for step in verification:
        p = text_frame.add_paragraph()
        p.text = step
        p.font.size = Pt(18)
        p.font.color.rgb = COLORS['text_dark']
        if step.startswith("   ") or step.startswith("Location:"):
            p.level = 1

def add_confidence_slide(prs):
    """Add confidence and limitations slide"""
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    
    title_shape = slide.shapes.title
    title_shape.text = "Confidence and Limitations"
    title_shape.text_frame.paragraphs[0].font.size = Pt(36)
    title_shape.text_frame.paragraphs[0].font.bold = True
    title_shape.text_frame.paragraphs[0].font.color.rgb = COLORS['primary']
    
    body_shape = slide.placeholders[1]
    text_frame = body_shape.text_frame
    text_frame.clear()
    
    # Confidence
    confidence = [
        "What We're Confident About:",
        "  • Revenue trajectory (5 years of consistent growth)",
        "  • Gross margins (stable at 87-90%)",
        "  • Balance sheet ($1.26B cash, negligible debt)",
        "  • SaaS model quality (subscription revenue)",
        "",
        "What We're Estimating:",
        "  • ARR and NRR (management metrics, not GAAP)",
        "  • AI monetization (uncertain, forward-looking)",
        "  • Competitive dynamics (unpredictable)",
        "",
        "What a Human Analyst Would Do Differently:",
        "  • Download and read full 10-K filings",
        "  • Attend earnings calls (live)",
        "  • Customer interviews",
        "  • Competitor product evaluation",
        "  • Management meetings",
    ]
    
    for item in confidence:
        p = text_frame.add_paragraph()
        p.text = item
        p.font.size = Pt(18)
        p.font.color.rgb = COLORS['text_dark']
        if item.startswith("  •") or item.startswith("  -"):
            p.level = 1

def add_recommendation_recap_slide(prs):
    """Add recommendation recap slide"""
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    
    title_shape = slide.shapes.title
    title_shape.text = "Recommendation Recap"
    title_shape.text_frame.paragraphs[0].font.size = Pt(36)
    title_shape.text_frame.paragraphs[0].font.bold = True
    title_shape.text_frame.paragraphs[0].font.color.rgb = COLORS['primary']
    
    body_shape = slide.placeholders[1]
    text_frame = body_shape.text_frame
    text_frame.clear()
    
    # Recommendation
    recommendation = [
        "RECOMMENDATION: BUY @ $42.00 (95% upside)",
        "",
        "Position Sizing: 2-3% of portfolio",
        "Holding Period: 12 months",
        "",
        "Catalysts:",
        "  • FY2027 Q1 earnings (June 2026)",
        "  • GAAP operating profitability (expected FY2028)",
        "  • AI feature expansion",
        "  • Potential M&A",
        "",
        "Key Risks:",
        "  • NRR decline below 110%",
        "  • GitHub Copilot execution",
        "  • Macro downturn",
    ]
    
    for item in recommendation:
        p = text_frame.add_paragraph()
        p.text = item
        p.font.size = Pt(20)
        p.font.color.rgb = COLORS['text_dark']
        if item.startswith("  •") or item.startswith("  -"):
            p.level = 1
        if item.startswith("RECOMMENDATION:"):
            p.font.bold = True
            p.font.color.rgb = COLORS['secondary']

def add_qa_slide(prs):
    """Add Q&A slide"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
    
    # Background
    background = slide.background
    background.fill.solid()
    background.fill.fore_color.rgb = COLORS['bg']
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(1), Inches(3), Inches(8), Inches(1))
    title_frame = title_box.text_frame
    title_frame.text = "Q&A"
    title_para = title_frame.paragraphs[0]
    title_para.font.size = Pt(72)
    title_para.font.bold = True
    title_para.font.color.rgb = COLORS['primary']
    title_para.alignment = PP_ALIGN.CENTER

# Add slides
add_title_slide(prs, "GitLab Inc. (GTLB)", "Board of Advisors Presentation | June 18, 2025")
add_recommendation_slide(prs, "BUY", "42.00", "21.51", "95% upside")
add_thesis_slide(prs, [
    "Large and Growing TAM ($30B+ DevOps market)",
    "Proven Growth Engine (30% CAGR revenue, 30% CAGR ARR)",
    "Clear Path to Profitability (operating margins improving 13pp annually)",
    "Strong Balance Sheet ($1.26B cash, negligible debt)",
    "What the Market Is Missing (AI monetization, operating leverage, ARR quality)",
])
add_financial_trajectory_slide(prs)
add_competitive_landscape_slide(prs)
add_valuation_slide(prs)
add_risk_assessment_slide(prs)
add_scenarios_slide(prs)
add_reasoning_trail_slide(prs)
add_dead_ends_slide(prs)
add_verification_slide(prs)
add_confidence_slide(prs)
add_recommendation_recap_slide(prs)
add_qa_slide(prs)

# Save presentation
output_path = '/workspace/deck/gtlb_board_presentation.pptx'
prs.save(output_path)
print(f"Presentation saved to: {output_path}")

# Export to PDF (requires PowerPoint or LibreOffice)
# This would typically be done manually or via command line:
# soffice --headless --convert-to pdf --outdir /workspace/deck /workspace/deck/gtlb_board_presentation.pptx
