#!/usr/bin/env python3
"""
Build the GitLab (GTLB) Board Presentation
Generates: /workspace/deck/gitlab_board_presentation.pptx

16 slides covering: recommendation, thesis, evidence, reasoning, verification, close.
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# ── Color Palette ──
BLUE = RGBColor(0x25, 0x63, 0xEB)
GREEN = RGBColor(0x05, 0x96, 0x69)
RED = RGBColor(0xDC, 0x26, 0x26)
PURPLE = RGBColor(0x7C, 0x3A, 0xED)
ORANGE = RGBColor(0xD9, 0x77, 0x06)
GRAY = RGBColor(0x64, 0x74, 0x8B)
DARK = RGBColor(0x1E, 0x29, 0x3B)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
BG = RGBColor(0xF8, 0xFA, 0xFC)
LIGHT_GRAY = RGBColor(0xE2, 0xE8, 0xF0)

# ── Helper Functions ──

def set_slide_bg(slide, color=BG):
    """Set slide background color."""
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb_color = color

def add_textbox(slide, left, top, width, height, text, font_size=18,
                bold=False, color=DARK, alignment=PP_ALIGN.LEFT, font_name='Calibri'):
    """Add a text box to a slide."""
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = font_name
    p.alignment = alignment
    return txBox

def add_multiline_textbox(slide, left, top, width, height, lines, font_size=14,
                          color=DARK, bold_first=False, spacing=1.2):
    """Add a text box with multiple lines."""
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, line in enumerate(lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = line
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.font.name = 'Calibri'
        p.space_after = Pt(font_size * spacing)
        if bold_first and i == 0:
            p.font.bold = True
    return txBox

def add_image(slide, left, top, width, height, path):
    """Add an image to a slide."""
    slide.shapes.add_picture(path, Inches(left), Inches(top), Inches(width), Inches(height))

def add_shape_box(slide, left, top, width, height, fill_color, text="", font_size=12, text_color=WHITE):
    """Add a colored shape box with optional text."""
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                    Inches(left), Inches(top), Inches(width), Inches(height))
    shape.fill.solid()
    shape.fill.fore_color.rgb_color = fill_color
    shape.line.fill.background()
    if text:
        tf = shape.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = text
        p.font.size = Pt(font_size)
        p.font.color.rgb = text_color
        p.font.name = 'Calibri'
        p.alignment = PP_ALIGN.CENTER
        tf.paragraphs[0].space_before = Pt(0)
        tf.paragraphs[0].space_after = Pt(0)
    return shape

def add_section_header(slide, text, top=0.3):
    """Add a section header bar."""
    # Background bar
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                    Inches(0), Inches(top), Inches(13.333), Inches(0.5))
    shape.fill.solid()
    shape.fill.fore_color.rgb_color = DARK
    shape.line.fill.background()
    # Text
    add_textbox(slide, 0.5, top + 0.05, 12, 0.4, text, font_size=16,
                bold=True, color=WHITE, alignment=PP_ALIGN.LEFT)

def add_footer(slide, text="Source: /input/repo/ — All data traceable via /audit/traces/"):
    """Add a footer to a slide."""
    add_textbox(slide, 0.5, 7.0, 12, 0.3, text, font_size=8,
                color=GRAY, alignment=PP_ALIGN.LEFT)

# ═══════════════════════════════════════════════════════════
# SLIDE 1: Title & Recommendation
# ═══════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank
set_slide_bg(slide, DARK)

# Title
add_textbox(slide, 1, 1.0, 11, 1.0, "GitLab Inc. (GTLB)", font_size=42,
            bold=True, color=WHITE, alignment=PP_ALIGN.LEFT)
add_textbox(slide, 1, 2.0, 11, 0.6, "Board of Advisors — Investment Recommendation", font_size=20,
            color=RGBColor(0x94, 0xA3, 0xB8), alignment=PP_ALIGN.LEFT)

# Recommendation box
add_shape_box(slide, 1, 3.2, 3.5, 1.8, GREEN, "BUY", font_size=36, text_color=WHITE)

# Price target
add_textbox(slide, 5, 3.2, 7, 0.5, "12-Month Price Target", font_size=14,
            color=RGBColor(0x94, 0xA3, 0xB8), alignment=PP_ALIGN.LEFT)
add_textbox(slide, 5, 3.7, 7, 1.0, "$42.00", font_size=48,
            bold=True, color=WHITE, alignment=PP_ALIGN.LEFT)
add_textbox(slide, 5, 4.7, 7, 0.5, "95% upside from current $21.51", font_size=16,
            color=GREEN, alignment=PP_ALIGN.LEFT)

# Weighted target
add_textbox(slide, 5, 5.5, 7, 0.5, "Probability-Weighted Target: $45.98 (114% upside)", font_size=14,
            color=RGBColor(0x94, 0xA3, 0xB8), alignment=PP_ALIGN.LEFT)

# Market cap
add_textbox(slide, 1, 5.5, 4, 0.5, "Market Cap: $3.66B | NYSE: GTLB", font_size=12,
            color=RGBColor(0x64, 0x74, 0x8B), alignment=PP_ALIGN.LEFT)

# Date
add_textbox(slide, 1, 6.5, 11, 0.3, "June 18, 2025 | Agent-Generated Analysis", font_size=10,
            color=RGBColor(0x47, 0x55, 0x69), alignment=PP_ALIGN.LEFT)

# ═══════════════════════════════════════════════════════════
# SLIDE 2: The Thesis
# ═══════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
add_section_header(slide, "THE THESIS — In the Agent's Own Words")

# 5 pillars
pillars = [
    ("1. Large & Growing TAM", "DevOps market estimated at $30B+, growing 15-20% annually. GitLab's full-stack approach positions it to capture share from fragmented point solutions.", BLUE),
    ("2. Proven Growth Engine", "Revenue grew from $424M (FY2023) to $955M (FY2026), a 30% CAGR. 115% NRR and $860M ARR drive predictable growth.", GREEN),
    ("3. Clear Path to Profitability", "Operating margins improved 13pp annually: -50% → -32% → -19% → -7%. GAAP profitability projected by FY2028.", PURPLE),
    ("4. Strong Balance Sheet", "$1.26B in cash & investments with negligible debt ($0.2M). Provides optionality for M&A, R&D, or weathering downturns.", ORANGE),
    ("5. What the Market Is Missing", "AI monetization potential, operating leverage trajectory, and ARR quality are not fully priced in at 2.56x EV/Revenue.", RED),
]

y_pos = 1.2
for title, desc, color in pillars:
    add_shape_box(slide, 0.5, y_pos, 0.4, 0.9, color)
    add_textbox(slide, 1.0, y_pos, 5, 0.35, title, font_size=14, bold=True, color=color)
    add_textbox(slide, 1.0, y_pos + 0.35, 11.5, 0.6, desc, font_size=11, color=DARK)
    y_pos += 1.1

add_footer(slide, "Source: /input/repo/memo/gitlab_investment_memo.md — Investment Thesis section")

# ═══════════════════════════════════════════════════════════
# SLIDE 3: What Would Change Our Mind
# ═══════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
add_section_header(slide, "WHAT WOULD CHANGE OUR MIND — Thesis Breakers")

add_textbox(slide, 0.5, 1.0, 12, 0.5, "Three scenarios that would invalidate the BUY recommendation:",
            font_size=14, color=GRAY)

breakers = [
    ("NRR Falls Below 110%",
     "NRR has declined from 120% (FY2022) to 115% (FY2026). If this trend continues below 110%, revenue growth would slow significantly, undermining the growth thesis. The 115% NRR still implies $990M in next-year revenue from existing customers — below 110%, this drops below $950M.",
     "What must be true: Enterprise customers stop expanding usage; AI features fail to drive upsells; competitive pressure forces price concessions.",
     RED),
    ("GitHub Captures Full DevOps Workflow",
     "Microsoft's GitHub Copilot has significant mindshare. If GitHub expands beyond code assistance into the full DevOps lifecycle (CI/CD, security, monitoring), GitLab's differentiation erodes. Currently, Copilot is a coding assistant, not a full platform.",
     "What must be true: Microsoft invests heavily in DevOps platform features; GitHub Actions becomes a full CI/CD replacement; enterprise customers consolidate on GitHub.",
     ORANGE),
    ("Macro Recession Cuts Enterprise Software Spend",
     "Enterprise software spending is cyclical. A recession could delay purchases, reduce NRR, and push back the profitability timeline. GitLab's $1.26B cash provides runway, but growth would slow.",
     "What must be true: GDP contraction >2%; enterprise IT budgets cut by 15%+; customer churn increases; new logo acquisition slows.",
     BLUE),
]

y_pos = 1.7
for title, desc, condition, color in breakers:
    add_shape_box(slide, 0.5, y_pos, 12.3, 1.6, LIGHT_GRAY)
    add_textbox(slide, 0.8, y_pos + 0.1, 11.5, 0.3, title, font_size=14, bold=True, color=color)
    add_textbox(slide, 0.8, y_pos + 0.45, 11.5, 0.5, desc, font_size=10, color=DARK)
    add_textbox(slide, 0.8, y_pos + 0.95, 11.5, 0.4, condition, font_size=9, color=GRAY)
    y_pos += 1.8

add_footer(slide, "Source: /input/repo/memo/gitlab_investment_memo.md — Risk Assessment section")

# ═══════════════════════════════════════════════════════════
# SLIDE 4: Financial Trajectory
# ═══════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
add_section_header(slide, "FINANCIAL TRAJECTORY — Revenue & Margins")

add_image(slide, 0.5, 1.0, 12.3, 5.5, "/workspace/assets/charts/01_revenue_margins.png")

# Key callouts
callouts = [
    "FY2024: FCF turns positive ($33M) — inflection point",
    "FY2026: Operating margin -7.4% — 13pp improvement per year",
    "FY2026: Revenue $955M — approaching $1B threshold for operating leverage",
]
y_pos = 6.6
for c in callouts:
    add_textbox(slide, 0.5, y_pos, 12, 0.25, "▸ " + c, font_size=10, color=DARK)
    y_pos += 0.25

add_footer(slide, "Source: /input/repo/extracted/income_statement_annual.csv, financial_summary.json")

# ═══════════════════════════════════════════════════════════
# SLIDE 5: SaaS Metrics
# ═══════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
add_section_header(slide, "SaaS METRICS — The Growth Engine")

add_image(slide, 0.5, 1.0, 12.3, 5.5, "/workspace/assets/charts/02_arr_nrr.png")

# Key metrics boxes
metrics = [
    ("ARR", "$860M", "FY2026"),
    ("NRR", "115%", "Declining from 120%"),
    ("Gross Margin", "87.4%", "Stable 87-90%"),
    ("Deferred Revenue", "$572M", "High visibility"),
]
x_pos = 0.5
for label, value, note in metrics:
    add_shape_box(slide, x_pos, 6.5, 2.8, 0.7, BLUE, f"{value}\n{label}", font_size=14)
    add_textbox(slide, x_pos, 7.15, 2.8, 0.2, note, font_size=8, color=GRAY, alignment=PP_ALIGN.CENTER)
    x_pos += 3.2

add_footer(slide, "Source: /input/repo/extracted/financial_summary.json, balance_sheet_annual.csv")

# ═══════════════════════════════════════════════════════════
# SLIDE 6: Competitive Position
# ═══════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
add_section_header(slide, "COMPETITIVE POSITION — Growth vs. Valuation")

add_image(slide, 0.5, 1.0, 12.3, 5.5, "/workspace/assets/charts/03_competitive_position.png")

# Key insight
add_textbox(slide, 0.5, 6.6, 12, 0.5, "▸ GitLab trades at 2.56x EV/Revenue vs. peer average of ~6.8x — a 62% discount for similar growth rates",
            font_size=11, bold=True, color=BLUE)
add_textbox(slide, 0.5, 7.0, 12, 0.3, "Comp set selected per ADR-002: 8 SaaS infrastructure/application peers with similar growth profiles",
            font_size=9, color=GRAY)

add_footer(slide, "Source: /input/repo/extracted/competitor_data.json")

# ═══════════════════════════════════════════════════════════
# SLIDE 7: The Mispricing Thesis
# ═══════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
add_section_header(slide, "THE MISPRICING THESIS — Why the Market Is Wrong")

# Three valuation anchors
anchors = [
    ("Current Price", "$21.51", "2.56x EV/Revenue", "Market's current assessment", RED),
    ("DCF Base Case", "$38.52", "10.5% WACC, 3% terminal", "Intrinsic value from cash flows", BLUE),
    ("Peer EV/Revenue", "$53.25", "7.03x peer average", "Relative valuation to peers", GREEN),
]

x_pos = 0.5
for title, value, basis, desc, color in anchors:
    add_shape_box(slide, x_pos, 1.5, 3.5, 2.5, color)
    add_textbox(slide, x_pos, 1.6, 3.5, 0.4, title, font_size=14, bold=True, color=WHITE, alignment=PP_ALIGN.CENTER)
    add_textbox(slide, x_pos, 2.1, 3.5, 0.8, value, font_size=36, bold=True, color=WHITE, alignment=PP_ALIGN.CENTER)
    add_textbox(slide, x_pos, 3.0, 3.5, 0.4, basis, font_size=10, color=WHITE, alignment=PP_ALIGN.CENTER)
    add_textbox(slide, x_pos, 3.5, 3.5, 0.4, desc, font_size=9, color=RGBColor(0xD0, 0xD0, 0xD0), alignment=PP_ALIGN.CENTER)
    x_pos += 4.0

# Our target
add_shape_box(slide, 5.5, 4.5, 2.3, 1.2, GREEN, "Our Target\n$42.00", font_size=18)
add_textbox(slide, 5.5, 5.7, 2.3, 0.4, "95% upside", font_size=12, bold=True, color=GREEN, alignment=PP_ALIGN.CENTER)

# Three factors
add_textbox(slide, 0.5, 6.0, 12, 0.3, "Three factors the market is underpricing:", font_size=13, bold=True, color=DARK)
factors = [
    "1. AI Monetization: GitLab Duo embedded in workflow creates natural upsell paths",
    "2. Operating Leverage: Fixed costs spread over growing revenue base past $1B",
    "3. Balance Sheet: $1.26B cash provides M&A optionality and downside protection",
]
y_pos = 6.3
for f in factors:
    add_textbox(slide, 0.5, y_pos, 12, 0.25, f, font_size=10, color=DARK)
    y_pos += 0.25

add_footer(slide, "Source: /input/repo/memo/gitlab_investment_memo.md — Valuation section")

# ═══════════════════════════════════════════════════════════
# SLIDE 8: Risk Assessment
# ═══════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
add_section_header(slide, "RISK ASSESSMENT — Prioritized by Probability & Impact")

add_image(slide, 0.5, 1.0, 12.3, 5.5, "/workspace/assets/charts/05_risk_heatmap.png")

add_textbox(slide, 0.5, 6.6, 12, 0.3, "▸ Highest priority: NRR decline (75% probability, 90% impact) — monitor quarterly NRR trends",
            font_size=10, color=DARK)
add_textbox(slide, 0.5, 6.9, 12, 0.3, "▸ Second priority: GitHub competition (50% probability, 85% impact) — watch Microsoft's DevOps roadmap",
            font_size=10, color=DARK)

add_footer(slide, "Source: /input/repo/memo/gitlab_investment_memo.md — Risk Assessment section")

# ═══════════════════════════════════════════════════════════
# SLIDE 9: Scenario Distribution
# ═══════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
add_section_header(slide, "BEAR / BASE / BULL — Probability-Weighted Distribution")

add_image(slide, 0.5, 1.0, 12.3, 5.5, "/workspace/assets/charts/04_scenario_distribution.png")

# Scenario details
scenarios_text = [
    ("Bear (25%)", "$18.58", "12% revenue growth, 12% WACC, 2% terminal growth", RED),
    ("Base (50%)", "$38.52", "20% revenue growth, 10.5% WACC, 3% terminal growth", BLUE),
    ("Bull (25%)", "$88.30", "28% revenue growth, 9% WACC, 4% terminal growth", GREEN),
]
x_pos = 0.5
for label, price, desc, color in scenarios_text:
    add_textbox(slide, x_pos, 6.6, 3.5, 0.25, label, font_size=11, bold=True, color=color)
    add_textbox(slide, x_pos, 6.85, 3.5, 0.2, price, font_size=14, bold=True, color=color)
    x_pos += 4.3

add_footer(slide, "Source: /input/repo/memo/gitlab_investment_memo.md — Scenario Analysis table")

# ═══════════════════════════════════════════════════════════
# SLIDE 10: Reasoning Trail
# ═══════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
add_section_header(slide, "THE REASONING TRAIL — How We Got Here")

add_image(slide, 0.3, 0.9, 12.7, 6.0, "/workspace/assets/charts/06_reasoning_trail.png")

add_footer(slide, "Source: /input/repo/decisions/ (ADR-001 through ADR-003), /input/repo/tool-log.md")

# ═══════════════════════════════════════════════════════════
# SLIDE 11: Dead Ends
# ═══════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
add_section_header(slide, "DEAD ENDS — What We Investigated and Rejected")

add_textbox(slide, 0.5, 1.0, 12, 0.5, "Honest accounting of what didn't work — boards trust analysts who show their work.",
            font_size=13, color=GRAY)

dead_ends = [
    ("SEC Filing Downloads", "Attempted 10-K/10-Q downloads from sec.gov. All returned 403 (Forbidden). Tried multiple User-Agent headers, session-based requests, curl, text versions. SEC actively blocks automated access.", "Used yfinance data instead (SEC-derived secondary source)"),
    ("GitLab IR Website", "DNS resolution failed for investors.gitlab.com. Could not access official earnings presentations or press releases.", "Relied on Seeking Alpha transcripts instead"),
    ("PRNewswire Press Releases", "Search URL returned 404. PRNewswire may have changed their search API or blocked automated access.", "No press releases included in analysis"),
    ("Initial CIK Lookup", "CIK 0001687226 mapped to 'Harper James P' not GitLab. Found correct CIK 0001653482 through SEC full-text search.", "Correct CIK used for all subsequent data retrieval"),
    ("Earnings Call Audio", "No accessible audio sources. Seeking Alpha 'Play Call' buttons require JavaScript execution.", "Relied on text transcripts only"),
]

y_pos = 1.7
for title, desc, resolution in dead_ends:
    add_shape_box(slide, 0.5, y_pos, 12.3, 0.95, LIGHT_GRAY)
    add_textbox(slide, 0.8, y_pos + 0.05, 5, 0.25, title, font_size=12, bold=True, color=RED)
    add_textbox(slide, 0.8, y_pos + 0.3, 11.5, 0.3, desc, font_size=9, color=DARK)
    add_textbox(slide, 0.8, y_pos + 0.6, 11.5, 0.25, "Resolution: " + resolution, font_size=9, color=GRAY)
    y_pos += 1.05

add_footer(slide, "Source: /input/repo/research/dead-ends.md")

# ═══════════════════════════════════════════════════════════
# SLIDE 12: Confidence & Limitations
# ═══════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
add_section_header(slide, "CONFIDENCE & LIMITATIONS — What We Know vs. What We Estimate")

add_image(slide, 0.5, 1.0, 12.3, 5.5, "/workspace/assets/charts/07_confidence_spectrum.png")

add_textbox(slide, 0.5, 6.6, 12, 0.3, "▸ We are most confident in revenue trajectory, gross margins, and balance sheet — these are facts, not estimates",
            font_size=10, color=DARK)
add_textbox(slide, 0.5, 6.9, 12, 0.3, "▸ We are least confident in AI monetization and competitive dynamics — these are forward-looking estimates",
            font_size=10, color=DARK)

add_footer(slide, "Source: /input/repo/memo/gitlab_investment_memo.md — Confidence and Limitations section")

# ═══════════════════════════════════════════════════════════
# SLIDE 13: How to Audit This Deck
# ═══════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
add_section_header(slide, "HOW TO AUDIT THIS DECK — Self-Verification Guide")

add_textbox(slide, 0.5, 1.0, 12, 0.5, "Any board member can verify any claim in this deck in under 2 minutes:",
            font_size=14, bold=True, color=DARK)

steps = [
    ("1. Clone the repo", "git clone <repo-url> && cd <repo-dir>", "Get the full audit trail"),
    ("2. Pick any number", "e.g., 'Revenue FY2026 = $955M' on Slide 4", "Choose any claim to verify"),
    ("3. Find the trace file", "cat audit/traces/trace_02_revenue_fy2026.md", "Trace file points to source"),
    ("4. Follow to source", "cat /input/repo/extracted/income_statement_annual.csv", "Source data in input repo"),
    ("5. Verify the number", "Total Revenue, 2026-01-31: 955224000.0 → $955.2M → $955M", "Number matches ✓"),
]

y_pos = 1.7
for title, command, desc in steps:
    add_shape_box(slide, 0.5, y_pos, 12.3, 0.85, LIGHT_GRAY)
    add_textbox(slide, 0.8, y_pos + 0.05, 4, 0.25, title, font_size=12, bold=True, color=BLUE)
    add_textbox(slide, 0.8, y_pos + 0.3, 11.5, 0.2, command, font_size=10, color=DARK)
    add_textbox(slide, 0.8, y_pos + 0.55, 11.5, 0.2, desc, font_size=9, color=GRAY)
    y_pos += 1.0

# Key directories
add_textbox(slide, 0.5, 6.8, 12, 0.5, "Key directories: /audit/traces/ (claim traces) | /audit/numbers.md (all numbers) | /audit/quotes.md (all quotes) | /assets/charts/ (reproducible scripts)",
            font_size=9, color=GRAY)

# ═══════════════════════════════════════════════════════════
# SLIDE 14: Design & Methodology Decisions
# ═══════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
add_section_header(slide, "DESIGN & METHODOLOGY — Every Choice Documented")

# ADR summary
add_textbox(slide, 0.5, 1.0, 6, 0.4, "Valuation Methodology (ADR-003)", font_size=13, bold=True, color=BLUE)
adr_items = [
    "Primary: DCF with FCF projections (GitLab has positive FCF, predictable revenue)",
    "Secondary: EV/Revenue peer comparison (standard SaaS metric)",
    "WACC: 10.5% (Rf=4.5%, Beta=1.2, ERP=5.5%)",
    "Terminal growth: 3.0% (conservative, below GDP growth)",
]
y_pos = 1.5
for item in adr_items:
    add_textbox(slide, 0.8, y_pos, 5.5, 0.25, "▸ " + item, font_size=10, color=DARK)
    y_pos += 0.3

add_textbox(slide, 0.5, 3.0, 6, 0.4, "Competitor Selection (ADR-002)", font_size=13, bold=True, color=GREEN)
comp_items = [
    "8 SaaS infrastructure/application peers with similar growth profiles",
    "Excluded GitHub (private, owned by MSFT) and ServiceNow (too large)",
    "Mix of profitable and growth-stage companies for range of multiples",
]
y_pos = 3.5
for item in comp_items:
    add_textbox(slide, 0.8, y_pos, 5.5, 0.25, "▸ " + item, font_size=10, color=DARK)
    y_pos += 0.3

add_textbox(slide, 7.5, 1.0, 5.5, 0.4, "Chart Conventions (ADR-005)", font_size=13, bold=True, color=PURPLE)
chart_items = [
    "Dual-axis: Revenue (bars) + margins (line) for financial trajectory",
    "Scatter: Growth vs. valuation for competitive positioning",
    "Distribution: Probability-weighted scenarios, not bullet points",
    "Heat map: Probability vs. impact for risk assessment",
    "Dependency graph: Top-down flow for reasoning trail",
]
y_pos = 1.5
for item in chart_items:
    add_textbox(slide, 7.8, y_pos, 5, 0.25, "▸ " + item, font_size=10, color=DARK)
    y_pos += 0.3

add_textbox(slide, 7.5, 3.0, 5.5, 0.4, "Color Palette (ADR-004)", font_size=13, bold=True, color=ORANGE)
color_items = [
    "Blue: Revenue, primary data, GitLab elements",
    "Green: Positive metrics, bull case, conclusion",
    "Red: Negative metrics, bear case, risk",
    "Purple: ARR, SaaS metrics, scenario analysis",
    "Orange: Warnings, ADRs, estimating category",
]
y_pos = 3.5
for item in color_items:
    add_textbox(slide, 7.8, y_pos, 5, 0.25, "▸ " + item, font_size=10, color=DARK)
    y_pos += 0.3

add_textbox(slide, 0.5, 6.5, 12, 0.5, "Every non-obvious aesthetic and analytical decision has an ADR in /decisions/. 'Used red for downside' is convention; 'chose sequential colormap over diverging because data has natural zero at current price' is an ADR.",
            font_size=10, color=GRAY)

add_footer(slide, "Source: /workspace/decisions/ (ADR-001 through ADR-005)")

# ═══════════════════════════════════════════════════════════
# SLIDE 15: Recommendation & Next Steps
# ═══════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, DARK)

add_textbox(slide, 1, 0.8, 11, 0.8, "RECOMMENDATION & NEXT STEPS", font_size=28,
            bold=True, color=WHITE, alignment=PP_ALIGN.LEFT)

# Recommendation
add_shape_box(slide, 1, 2.0, 5, 1.5, GREEN, "BUY\n$42.00 | 95% Upside", font_size=24)

# Position sizing
add_textbox(slide, 7, 2.0, 5, 0.3, "Position Sizing", font_size=14, bold=True, color=WHITE)
add_textbox(slide, 7, 2.4, 5, 0.5, "2-3% of portfolio", font_size=20, bold=True, color=GREEN)
add_textbox(slide, 7, 2.9, 5, 0.3, "Moderate position given high conviction\nbut timing uncertainty", font_size=10, color=RGBColor(0x94, 0xA3, 0xB8))

# Holding period
add_textbox(slide, 7, 3.7, 5, 0.3, "Holding Period", font_size=14, bold=True, color=WHITE)
add_textbox(slide, 7, 4.1, 5, 0.5, "12 months", font_size=20, bold=True, color=GREEN)
add_textbox(slide, 7, 4.6, 5, 0.3, "Re-evaluate at FY2027 Q1 earnings", font_size=10, color=RGBColor(0x94, 0xA3, 0xB8))

# Catalysts
add_textbox(slide, 1, 4.0, 5, 0.3, "Key Catalysts", font_size=14, bold=True, color=WHITE)
catalysts = [
    "FY2027 Q1 earnings (June 2026): First quarter with AI revenue contribution",
    "GAAP operating profitability: Expected FY2028",
    "AI feature expansion: New AI capabilities driving upsell",
    "Potential M&A: $1.26B cash enables strategic acquisitions",
]
y_pos = 4.4
for c in catalysts:
    add_textbox(slide, 1.2, y_pos, 4.5, 0.25, "▸ " + c, font_size=10, color=RGBColor(0x94, 0xA3, 0xB8))
    y_pos += 0.3

# Summary
add_textbox(slide, 1, 6.2, 11, 0.5, "GitLab is a high-quality SaaS company with a clear path to profitability, strong SaaS metrics, and a differentiated competitive position. The current valuation of 2.6x EV/Revenue is well below peer averages.",
            font_size=11, color=RGBColor(0x94, 0xA3, 0xB8))

add_textbox(slide, 1, 6.8, 11, 0.3, "Source: /input/repo/memo/gitlab_investment_memo.md — Recommendation section",
            font_size=8, color=RGBColor(0x47, 0x55, 0x69))

# ═══════════════════════════════════════════════════════════
# SLIDE 16: Q&A / Appendix
# ═══════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, DARK)

add_textbox(slide, 1, 1.0, 11, 0.8, "Q&A", font_size=42,
            bold=True, color=WHITE, alignment=PP_ALIGN.LEFT)

# Key numbers at a glance
add_textbox(slide, 1, 2.5, 5, 0.3, "Key Numbers at a Glance", font_size=14, bold=True, color=WHITE)
key_nums = [
    ("Current Price", "$21.51"),
    ("Price Target", "$42.00"),
    ("Upside", "95%"),
    ("Market Cap", "$3.66B"),
    ("Revenue (FY2026)", "$955M"),
    ("ARR (FY2026)", "$860M"),
    ("NRR", "115%"),
    ("Cash & Investments", "$1.26B"),
    ("EV/Revenue", "2.56x"),
    ("Peer Avg EV/Revenue", "~6.8x"),
]
x_pos = 1.0
y_pos = 3.0
for label, value in key_nums:
    add_textbox(slide, x_pos, y_pos, 2.2, 0.2, label, font_size=10, color=RGBColor(0x94, 0xA3, 0xB8))
    add_textbox(slide, x_pos, y_pos + 0.2, 2.2, 0.25, value, font_size=12, bold=True, color=WHITE)
    x_pos += 2.5
    if x_pos > 8:
        x_pos = 1.0
        y_pos += 0.6

# Repo info
add_textbox(slide, 1, 6.0, 11, 0.3, "Full audit trail: /audit/ | Chart scripts: /assets/charts/ | Decision records: /decisions/",
            font_size=10, color=RGBColor(0x64, 0x74, 0x8B))
add_textbox(slide, 1, 6.4, 11, 0.3, "Disclaimer: This memo is for informational purposes only and does not constitute investment advice.",
            font_size=9, color=RGBColor(0x47, 0x55, 0x69))

# ═══════════════════════════════════════════════════════════
# SAVE
# ═══════════════════════════════════════════════════════════
output_path = "/workspace/deck/gitlab_board_presentation.pptx"
prs.save(output_path)
print(f"Presentation saved: {output_path}")
print(f"Total slides: {len(prs.slides)}")
