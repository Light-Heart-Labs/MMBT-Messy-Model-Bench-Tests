#!/usr/bin/env python3
"""
Build the Board of Advisors Presentation for GitLab (GTLB)
15 slides, designed for 30-minute presentation + 15-minute Q&A.

Output: /workspace/deck/gitlab_board_presentation.pptx
"""

import os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# Colors
TEAL = RGBColor(0x0D, 0x94, 0x88)
TEAL_LIGHT = RGBColor(0x99, 0xF6, 0xE4)
BLUE = RGBColor(0x25, 0x63, 0xEB)
RED = RGBColor(0xDC, 0x26, 0x26)
GREEN = RGBColor(0x16, 0xA3, 0x4A)
GRAY = RGBColor(0x6B, 0x72, 0x80)
LIGHT_GRAY = RGBColor(0xF3, 0xF4, 0xF6)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
BLACK = RGBColor(0x11, 0x11, 0x11)
DARK_GRAY = RGBColor(0x37, 0x41, 0x51)
ORANGE = RGBColor(0xF5, 0x9E, 0x0B)
PURPLE = RGBColor(0x8B, 0x5C, 0xF6)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

def add_background(slide, color=WHITE):
    """Set slide background color."""
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_text_box(slide, left, top, width, height, text, font_size=14, 
                 font_color=BLACK, bold=False, alignment=PP_ALIGN.LEFT,
                 font_name='Calibri'):
    """Add a text box to the slide."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = font_color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    return txBox

def add_bullet_list(slide, left, top, width, height, items, font_size=14,
                    font_color=BLACK, bold_first=False, font_name='Calibri'):
    """Add a bulleted list."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = item
        p.font.size = Pt(font_size)
        p.font.color.rgb = font_color
        p.font.name = font_name
        p.level = 0
        p.space_after = Pt(6)
        if bold_first and i == 0:
            p.font.bold = True
    return txBox

def add_image(slide, left, top, width, height, path):
    """Add an image to the slide."""
    if os.path.exists(path):
        slide.shapes.add_picture(path, left, top, width, height)
        return True
    return False

def add_shape_rect(slide, left, top, width, height, fill_color, border_color=None):
    """Add a rounded rectangle shape."""
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if border_color:
        shape.line.color.rgb = border_color
        shape.line.width = Pt(1)
    else:
        shape.line.fill.background()
    return shape

def add_footer(slide, text="Source: /input/repo/ | Audit trail: /audit/traces/"):
    """Add a footer to the slide."""
    add_text_box(slide, Inches(0.5), Inches(7.0), Inches(12), Inches(0.4),
                 text, font_size=8, font_color=GRAY, alignment=PP_ALIGN.LEFT)

# ============================================================
# SLIDE 1: Title & Recommendation
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank
add_background(slide, WHITE)

# Title area
add_shape_rect(slide, Inches(0), Inches(0), Inches(13.333), Inches(3.5), TEAL)

add_text_box(slide, Inches(1), Inches(0.5), Inches(11), Inches(1),
             "GitLab Inc. (GTLB)", font_size=40, font_color=WHITE, bold=True)
add_text_box(slide, Inches(1), Inches(1.4), Inches(11), Inches(0.8),
             "Board of Advisors Presentation — Investment Analysis", font_size=20, font_color=TEAL_LIGHT)

# Recommendation box
add_shape_rect(slide, Inches(1), Inches(2.3), Inches(4), Inches(1), WHITE)
add_text_box(slide, Inches(1.2), Inches(2.35), Inches(3.6), Inches(0.5),
             "RECOMMENDATION", font_size=14, font_color=TEAL, bold=True, alignment=PP_ALIGN.CENTER)
add_text_box(slide, Inches(1.2), Inches(2.8), Inches(3.6), Inches(0.5),
             "BUY", font_size=36, font_color=TEAL, bold=True, alignment=PP_ALIGN.CENTER)

# Price target box
add_shape_rect(slide, Inches(5.5), Inches(2.3), Inches(4), Inches(1), WHITE)
add_text_box(slide, Inches(5.7), Inches(2.35), Inches(3.6), Inches(0.5),
             "12-MONTH PRICE TARGET", font_size=14, font_color=BLUE, bold=True, alignment=PP_ALIGN.CENTER)
add_text_box(slide, Inches(5.7), Inches(2.8), Inches(3.6), Inches(0.5),
             "$42.00", font_size=36, font_color=BLUE, bold=True, alignment=PP_ALIGN.CENTER)

# Upside box
add_shape_rect(slide, Inches(10), Inches(2.3), Inches(3), Inches(1), WHITE)
add_text_box(slide, Inches(10.2), Inches(2.35), Inches(2.6), Inches(0.5),
             "UPSIDE", font_size=14, font_color=GREEN, bold=True, alignment=PP_ALIGN.CENTER)
add_text_box(slide, Inches(10.2), Inches(2.8), Inches(2.6), Inches(0.5),
             "+95%", font_size=36, font_color=GREEN, bold=True, alignment=PP_ALIGN.CENTER)

# Key stats
add_text_box(slide, Inches(1), Inches(3.7), Inches(11), Inches(0.4),
             "Current Price: $21.51  |  Market Cap: $3.66B  |  Sector: Software — Infrastructure  |  Shares Outstanding: 170.1M",
             font_size=14, font_color=DARK_GRAY)

# Confidence meter
add_text_box(slide, Inches(1), Inches(4.3), Inches(11), Inches(0.4),
             "Confidence Level: 7/10 — Strong conviction in business model, moderate uncertainty on timing of profitability",
             font_size=13, font_color=GRAY)

# Confidence bar
add_shape_rect(slide, Inches(1), Inches(4.7), Inches(11), Inches(0.3), LIGHT_GRAY)
add_shape_rect(slide, Inches(1), Inches(4.7), Inches(7.7), Inches(0.3), TEAL)
add_text_box(slide, Inches(1), Inches(5.1), Inches(11), Inches(0.3),
             "0          2          4          6          8          10",
             font_size=9, font_color=GRAY, alignment=PP_ALIGN.CENTER)

# Image
add_image(slide, Inches(8.5), Inches(5.5), Inches(4.5), Inches(1.5),
          "/workspace/assets/charts/10-stock-price.png")

add_footer(slide, "Source: /input/repo/memo/gitlab_investment_memo.md, /input/repo/extracted/company_info.json | Trace: /audit/traces/slide1_recommendation.md")

# ============================================================
# SLIDE 2: The Thesis in One Slide
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_background(slide, WHITE)

add_text_box(slide, Inches(0.8), Inches(0.3), Inches(11), Inches(0.6),
             "The Investment Thesis — In the Agent's Own Words",
             font_size=28, font_color=TEAL, bold=True)

add_text_box(slide, Inches(0.8), Inches(0.9), Inches(11), Inches(0.4),
             "GitLab is a high-quality SaaS company with a clear path to profitability, strong SaaS metrics, and a differentiated competitive position.",
             font_size=14, font_color=GRAY)

thesis_items = [
    "1. Large and Growing TAM — DevOps and developer tools market estimated at $30B+, growing 15-20% annually. GitLab's full-stack approach positions it to capture share from fragmented point solutions.",
    "2. Proven Growth Engine — Revenue grew from $424M (FY2023) to $955M (FY2026), 30% CAGR. ARR grew from $400M to $860M. NRR of 115% means existing customers expand 15% annually.",
    "3. Clear Path to Profitability — Operating margins improved from -49.8% (FY2023) to -7.4% (FY2026). FCF turned positive in FY2024 ($33M) and reached $222M in FY2026. GAAP profitability projected by FY2028.",
    "4. Strong Balance Sheet — $1.26B in cash/investments with negligible debt ($0.2M). Provides optionality for M&A, R&D, or weathering downturns.",
    "5. What the Market Is Missing — AI monetization potential (GitLab Duo), operating leverage as revenue passes $1B, and ARR quality ($860M with 115% NRR = ~$990M next-year revenue from existing customers)."
]

add_bullet_list(slide, Inches(0.8), Inches(1.5), Inches(11.5), Inches(5.0),
                thesis_items, font_size=13, font_color=DARK_GRAY)

add_footer(slide, "Source: /input/repo/memo/gitlab_investment_memo.md (Investment Thesis section) | Trace: /audit/traces/slide2_revenue_growth.md")

# ============================================================
# SLIDE 3: What Would Change the Recommendation
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_background(slide, WHITE)

add_text_box(slide, Inches(0.8), Inches(0.3), Inches(11), Inches(0.6),
             "What Would Change the Recommendation",
             font_size=28, font_color=RED, bold=True)

add_text_box(slide, Inches(0.8), Inches(0.9), Inches(11), Inches(0.4),
             "A good analyst shows the exit before showing the entry. Here are the 3 conditions that would flip BUY to HOLD or SELL.",
             font_size=14, font_color=GRAY)

# Three condition cards
conditions = [
    ("NRR Falls Below 110%", RED, 
     "If NRR declines from 115% to below 110%, revenue growth would slow significantly. This would indicate customers are not expanding usage, which undermines the core growth engine.\n\nWhat would have to be true: GitHub Copilot or Atlassian AI offers compelling alternatives that reduce GitLab's expansion revenue."),
    ("GitHub Copilot Eats Share", ORANGE,
     "If Microsoft's GitHub Copilot expands beyond coding assistance into full DevOps functionality, GitLab's open-core distribution advantage could erode.\n\nWhat would have to be true: GitHub integrates full CI/CD, security scanning, and monitoring capabilities, and developers prefer the GitHub ecosystem over GitLab's."),
    ("Macro Recession Hits Enterprise Software", BLUE,
     "A recession could delay enterprise software purchases, reduce NRR, and extend the path to profitability.\n\nWhat would have to be true: Enterprise software spending declines by 15%+ for 2+ consecutive quarters, and GitLab's pipeline dries up.")
]

for i, (title, color, desc) in enumerate(conditions):
    left = Inches(0.5 + i * 4.2)
    add_shape_rect(slide, left, Inches(1.5), Inches(3.9), Inches(0.6), color)
    add_text_box(slide, left, Inches(1.52), Inches(3.9), Inches(0.55),
                 title, font_size=14, font_color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
    add_text_box(slide, left, Inches(2.2), Inches(3.9), Inches(4.5),
                 desc, font_size=11, font_color=DARK_GRAY)

add_footer(slide, "Source: /input/repo/memo/gitlab_investment_memo.md (Risk Assessment section) | Trace: /audit/traces/slide3_nrr_threshold.md")

# ============================================================
# SLIDE 4: Financial Trajectory
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_background(slide, WHITE)

add_text_box(slide, Inches(0.8), Inches(0.3), Inches(11), Inches(0.6),
             "Financial Trajectory: Revenue, Margin & FCF",
             font_size=28, font_color=TEAL, bold=True)

add_text_box(slide, Inches(0.8), Inches(0.9), Inches(11), Inches(0.4),
             "Revenue growing 30% CAGR, margins improving rapidly, FCF turned positive in FY2024",
             font_size=14, font_color=GRAY)

# Chart
add_image(slide, Inches(0.5), Inches(1.4), Inches(8.5), Inches(5.0),
          "/workspace/assets/charts/01-financial-trajectory.png")

# Key metrics sidebar
metrics = [
    "Revenue CAGR (FY2023-FY2026): 30%",
    "Operating Margin Improvement: 42 pp",
    "FCF Turned Positive: FY2024",
    "Gross Margin: Stable 87-90%",
    "ARR Growth: 33% → 26% (moderating)",
    "NRR: 118% → 115% (declining but strong)"
]

add_bullet_list(slide, Inches(9.3), Inches(1.5), Inches(3.5), Inches(4.5),
                metrics, font_size=12, font_color=DARK_GRAY)

add_footer(slide, "Source: /input/repo/extracted/income_statement_annual.csv, cash_flow_annual.csv, financial_summary.json | Trace: /audit/traces/slide4_revenue.md")

# ============================================================
# SLIDE 5: Competitive Position
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_background(slide, WHITE)

add_text_box(slide, Inches(0.8), Inches(0.3), Inches(11), Inches(0.6),
             "Competitive Position: Growth vs. Margin",
             font_size=28, font_color=TEAL, bold=True)

add_text_box(slide, Inches(0.8), Inches(0.9), Inches(11), Inches(0.4),
             "GitLab is the only large-cap SaaS at 2.6x EV/Revenue with 23% growth — the market is pricing it like a low-growth company",
             font_size=14, font_color=GRAY)

# Chart
add_image(slide, Inches(0.5), Inches(1.4), Inches(8.5), Inches(5.0),
          "/workspace/assets/charts/02-competitive-position.png")

# Why these comps
add_text_box(slide, Inches(9.3), Inches(1.5), Inches(3.5), Inches(0.4),
             "Why These Comps", font_size=14, font_color=TEAL, bold=True)

comp_reasons = [
    "• SaaS business model",
    "• Similar growth profiles",
    "• Publicly traded",
    "• Mix of profitable",
    "  and growth-stage",
    "",
    "Excluded:",
    "• GitHub (private)",
    "• ServiceNow (too large)",
    "• Palo Alto (different"
]

add_bullet_list(slide, Inches(9.3), Inches(2.0), Inches(3.5), Inches(4.5),
                comp_reasons, font_size=11, font_color=DARK_GRAY)

add_footer(slide, "Source: /input/repo/extracted/competitor_data.json | Trace: /audit/traces/slide5_gtlb_ev_rev.md")

# ============================================================
# SLIDE 6: The Mispricing Thesis
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_background(slide, WHITE)

add_text_box(slide, Inches(0.8), Inches(0.3), Inches(11), Inches(0.6),
             "The Mispricing Thesis: Why the Market Is Wrong",
             font_size=28, font_color=TEAL, bold=True)

# Three columns
col_data = [
    ("Sell-Side Consensus", "$30.79", "24 analysts", 
     "Most focus on near-term profitability timeline. Limited analysis of AI monetization potential. NRR decline from 120% to 115% seen as negative, but still strong.",
     RED),
    ("Our DCF Base Case", "$38.52", "79% upside",
     "DCF with WACC 10.5%, terminal growth 3%. Captures the value of future cash generation as margins improve and revenue scales.",
     BLUE),
    ("Our 12-Month Target", "$42.00", "95% upside",
     "Midpoint between DCF ($38.52) and probability-weighted target ($45.98). Represents 3.5x FY2027E EV/Revenue, below peer average of 7.0x.",
     GREEN),
]

for i, (title, price, upside, desc, color) in enumerate(col_data):
    left = Inches(0.5 + i * 4.2)
    add_shape_rect(slide, left, Inches(1.3), Inches(3.9), Inches(0.8), color)
    add_text_box(slide, left, Inches(1.32), Inches(3.9), Inches(0.4),
                 title, font_size=14, font_color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
    add_text_box(slide, left, Inches(1.7), Inches(3.9), Inches(0.4),
                 f"{price} ({upside})", font_size=20, font_color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
    add_text_box(slide, left, Inches(2.3), Inches(3.9), Inches(3.5),
                 desc, font_size=11, font_color=DARK_GRAY)

# Key mispricing factors
add_text_box(slide, Inches(0.8), Inches(5.5), Inches(11), Inches(0.4),
             "Key Mispricing Factors:", font_size=16, font_color=TEAL, bold=True)

mispricing = [
    "1. AI Monetization — GitLab Duo embedded in developer workflow, creating natural upsell paths the market hasn't priced in",
    "2. Operating Leverage — As revenue grows past $1B, fixed costs spread over larger base. Projected operating margins: -7% (FY2026) → +5% (FY2029) → +12% (FY2031)",
    "3. Balance Sheet Optionality — $1.26B cash/investments with $0.2M debt provides M&A, R&D, or downturn optionality"
]

add_bullet_list(slide, Inches(0.8), Inches(5.9), Inches(11.5), Inches(1.5),
                mispricing, font_size=12, font_color=DARK_GRAY)

add_footer(slide, "Source: /input/repo/extracted/company_info.json, /input/repo/memo/gitlab_investment_memo.md | Trace: /audit/traces/slide6_consensus.md")

# ============================================================
# SLIDE 7: Risk Assessment
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_background(slide, WHITE)

add_text_box(slide, Inches(0.8), Inches(0.3), Inches(11), Inches(0.6),
             "Risk Assessment: Prioritized by Probability & Impact",
             font_size=28, font_color=RED, bold=True)

# Chart
add_image(slide, Inches(0.5), Inches(1.2), Inches(7.5), Inches(4.5),
          "/workspace/assets/charts/05-risk-matrix.png")

# Risk details
add_text_box(slide, Inches(8.3), Inches(1.2), Inches(4.5), Inches(0.4),
             "Top Risks", font_size=16, font_color=RED, bold=True)

risk_details = [
    "1. NRR Decline (High/High)",
    "   NRR fell 120%→115%. If <110%, growth slows.",
    "",
    "2. GitHub Competition (Med/High)",
    "   Copilot expanding beyond coding.",
    "",
    "3. Macro Downturn (Med/High)",
    "   Enterprise software spending cuts.",
    "",
    "4. AI Monetization Failure (Med/Med)",
    "   AI features don't drive upsells.",
    "",
    "5. Open-Source Cannibalization (Low/Med)",
    "   Self-hosted users prefer free version.",
    "",
    "6. Key Person Risk (Low/Low)",
    "   CEO Sid Sijbrandij is key visionary."
]

add_bullet_list(slide, Inches(8.3), Inches(1.7), Inches(4.5), Inches(5.0),
                risk_details, font_size=11, font_color=DARK_GRAY)

add_footer(slide, "Source: /input/repo/memo/gitlab_investment_memo.md (Risk Assessment) | Trace: /audit/traces/slide7_nrr_decline.md")

# ============================================================
# SLIDE 8: Bear/Base/Bull Scenarios
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_background(slide, WHITE)

add_text_box(slide, Inches(0.8), Inches(0.3), Inches(11), Inches(0.6),
             "Scenario Analysis: Probability-Weighted Price Distribution",
             font_size=28, font_color=TEAL, bold=True)

add_text_box(slide, Inches(0.8), Inches(0.9), Inches(11), Inches(0.4),
             "The price target is a distribution, not a number. Expected value: $45.98 | 12-month target: $42.00",
             font_size=14, font_color=GRAY)

# Chart
add_image(slide, Inches(0.5), Inches(1.4), Inches(9.5), Inches(5.0),
          "/workspace/assets/charts/03-scenario-distribution.png")

# Scenario table
add_text_box(slide, Inches(10.3), Inches(1.5), Inches(2.8), Inches(0.4),
             "Scenario Summary", font_size=14, font_color=TEAL, bold=True)

scenario_table = [
    "Bear (25%): $18.58",
    "  -14% from current",
    "  12% revenue growth",
    "  12.0% WACC",
    "",
    "Base (50%): $38.52",
    "  +79% from current",
    "  20% revenue growth",
    "  10.5% WACC",
    "",
    "Bull (25%): $88.30",
    "  +310% from current",
    "  28% revenue growth",
    "  9.0% WACC",
    "",
    "Expected: $45.98",
    "  +114% from current"
]

add_bullet_list(slide, Inches(10.3), Inches(2.0), Inches(2.8), Inches(5.0),
                scenario_table, font_size=10, font_color=DARK_GRAY)

add_footer(slide, "Source: /input/repo/memo/gitlab_investment_memo.md (Scenario Analysis) | Trace: /audit/traces/slide8_scenarios.md")

# ============================================================
# SLIDE 9: DCF Bridge
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_background(slide, WHITE)

add_text_box(slide, Inches(0.8), Inches(0.3), Inches(11), Inches(0.6),
             "DCF Bridge: From $21.51 to $42.00",
             font_size=28, font_color=TEAL, bold=True)

add_text_box(slide, Inches(0.8), Inches(0.9), Inches(11), Inches(0.4),
             "How we get from the current price to the target — each component is traceable to the input repo",
             font_size=14, font_color=GRAY)

# Chart
add_image(slide, Inches(0.5), Inches(1.4), Inches(9.5), Inches(5.0),
          "/workspace/assets/charts/09-dcf-bridge.png")

# Key assumptions
add_text_box(slide, Inches(10.3), Inches(1.5), Inches(2.8), Inches(0.4),
             "Key Assumptions", font_size=14, font_color=TEAL, bold=True)

assumptions = [
    "WACC: 10.5%",
    "  Rf: 4.5% (10Y Treasury)",
    "  Beta: 1.2 (SaaS sector)",
    "  ERP: 5.5%",
    "",
    "Terminal Growth: 3.0%",
    "  Conservative, below GDP",
    "",
    "Revenue CAGR: 20%",
    "  FY2027-FY2031",
    "",
    "FCF Margin: 25-28%",
    "  By FY2031",
    "",
    "Shares: 170.1M",
    "  No buybacks assumed"
]

add_bullet_list(slide, Inches(10.3), Inches(2.0), Inches(2.8), Inches(5.0),
                assumptions, font_size=10, font_color=DARK_GRAY)

add_footer(slide, "Source: /input/repo/decisions/003-valuation-methodology.md, /input/repo/memo/gitlab_investment_memo.md | Trace: /audit/traces/slide9_wacc.md")

# ============================================================
# SLIDE 10: The Reasoning Trail
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_background(slide, WHITE)

add_text_box(slide, Inches(0.8), Inches(0.3), Inches(11), Inches(0.6),
             "The Reasoning Trail: How the Agent Moved from Data to Conclusion",
             font_size=26, font_color=TEAL, bold=True)

add_text_box(slide, Inches(0.8), Inches(0.9), Inches(11), Inches(0.4),
             "This is the capability demonstration. Every node is a decision point or data artifact. Every edge is a data flow.",
             font_size=14, font_color=GRAY)

# Chart
add_image(slide, Inches(0.3), Inches(1.3), Inches(12.7), Inches(5.5),
          "/workspace/assets/charts/06-reasoning-trail.png")

add_footer(slide, "Source: /input/repo/decisions/, /input/repo/research/, /input/repo/memo/ | Trace: /audit/traces/slide10_transcripts.md")

# ============================================================
# SLIDE 11: Dead Ends
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_background(slide, WHITE)

add_text_box(slide, Inches(0.8), Inches(0.3), Inches(11), Inches(0.6),
             "Dead Ends: Hypotheses Investigated and Rejected",
             font_size=28, font_color=ORANGE, bold=True)

add_text_box(slide, Inches(0.8), Inches(0.9), Inches(11), Inches(0.4),
             "Boards trust analysts who show their work. Here are the 5 approaches that failed and what we learned.",
             font_size=14, font_color=GRAY)

# Chart
add_image(slide, Inches(0.5), Inches(1.4), Inches(9.5), Inches(4.5),
          "/workspace/assets/charts/07-dead-ends.png")

# Lessons learned
add_text_box(slide, Inches(10.3), Inches(1.5), Inches(2.8), Inches(0.4),
             "Lessons Learned", font_size=14, font_color=TEAL, bold=True)

lessons = [
    "SEC blocks automated",
    "  access — use yfinance",
    "  as proxy for filings",
    "",
    "CIK databases can be",
    "  outdated — verify",
    "  through SEC systems",
    "",
    "Text transcripts are",
    "  sufficient for analysis",
    "  when audio is blocked",
    "",
    "Document every dead end",
    "  — it proves rigor and",
    "  prevents repetition"
]

add_bullet_list(slide, Inches(10.3), Inches(2.0), Inches(2.8), Inches(4.5),
                lessons, font_size=11, font_color=DARK_GRAY)

add_footer(slide, "Source: /input/repo/research/dead-ends.md | Trace: /audit/traces/slide11_dead_ends.md")

# ============================================================
# SLIDE 12: Confidence & Limitations
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_background(slide, WHITE)

add_text_box(slide, Inches(0.8), Inches(0.3), Inches(11), Inches(0.6),
             "Confidence & Limitations: What We Know vs. What We Estimate",
             font_size=26, font_color=TEAL, bold=True)

# Chart
add_image(slide, Inches(0.3), Inches(1.0), Inches(6.5), Inches(5.5),
          "/workspace/assets/charts/08-confidence-limitations.png")

# What we're confident about
add_text_box(slide, Inches(7.2), Inches(1.0), Inches(5.5), Inches(0.4),
             "What We're Confident About", font_size=16, font_color=TEAL, bold=True)

confident_items = [
    "✓ Revenue trajectory — 5 years of consistent",
    "  growth data provides strong basis",
    "",
    "✓ Gross margins — Stable at 87-90%,",
    "  well-documented across filings",
    "",
    "✓ Balance sheet — $1.26B cash/investments",
    "  with negligible debt is a fact, not estimate",
    "",
    "✓ SaaS model quality — Subscription revenue",
    "  provides high visibility and predictability"
]

add_bullet_list(slide, Inches(7.2), Inches(1.5), Inches(5.5), Inches(3.0),
                confident_items, font_size=12, font_color=DARK_GRAY)

# What we're estimating
add_text_box(slide, Inches(7.2), Inches(4.5), Inches(5.5), Inches(0.4),
             "What We're Estimating", font_size=16, font_color=RED, bold=True)

estimating_items = [
    "⚠ ARR and NRR — Management metrics,",
    "  not GAAP figures. Derived from earnings",
    "  call transcripts, may differ from official.",
    "",
    "⚠ AI monetization — Revenue impact of AI",
    "  features is uncertain and forward-looking.",
    "",
    "⚠ Competitive dynamics — Pace of GitHub/",
    "  Atlassian AI feature development is",
    "  unpredictable."
]

add_bullet_list(slide, Inches(7.2), Inches(5.0), Inches(5.5), Inches(2.5),
                estimating_items, font_size=12, font_color=DARK_GRAY)

add_footer(slide, "Source: /input/repo/memo/gitlab_investment_memo.md (Confidence and Limitations) | Trace: /audit/traces/slide12_revenue_confidence.md")

# ============================================================
# SLIDE 13: How to Audit This Deck
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_background(slide, WHITE)

add_text_box(slide, Inches(0.8), Inches(0.3), Inches(11), Inches(0.6),
             "How to Audit This Deck — The Self-Audit Slide",
             font_size=28, font_color=TEAL, bold=True)

add_text_box(slide, Inches(0.8), Inches(0.9), Inches(11), Inches(0.4),
             "This is the most important slide for a skeptical technical board. It shows how to verify any claim in under 2 minutes.",
             font_size=14, font_color=GRAY)

# Step-by-step
steps = [
    ("1. Clone the repo", "git clone /workspace/ — This repo contains everything: the deck, the charts, the audit trail, and the source data."),
    ("2. Pick any claim", "Find a number on any slide. E.g., Slide 4 says 'Revenue: $955M (FY2026)'."),
    ("3. Follow the trace", "Go to /audit/traces/slide4_revenue.md. It points to /input/repo/extracted/income_statement_annual.csv, row 'Total Revenue', column '2026-01-31'."),
    ("4. Verify the number", "Open the CSV. The value is 955,224,000. Rounded to millions = $955M. Match confirmed."),
    ("5. Check the chart script", "Go to /assets/charts/01-financial-trajectory.py. Run it. You get back the exact chart from the deck."),
    ("6. Check the quotes", "Go to /audit/quotes.md. It shows the full paragraph context, not just the quote. Verify the quote isn't ripped out of context.")
]

for i, (step, desc) in enumerate(steps):
    top = Inches(1.5 + i * 0.95)
    add_shape_rect(slide, Inches(0.5), top, Inches(2.5), Inches(0.7), TEAL)
    add_text_box(slide, Inches(0.6), top, Inches(2.3), Inches(0.65),
                 step, font_size=13, font_color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
    add_text_box(slide, Inches(3.2), top, Inches(9.5), Inches(0.7),
                 desc, font_size=12, font_color=DARK_GRAY)

# Key files
add_text_box(slide, Inches(0.8), Inches(7.0), Inches(11), Inches(0.3),
             "Key files: /audit/traces/ (claim traces) | /audit/numbers.md (all numbers) | /audit/quotes.md (all quotes) | /audit/reconciliation.md (spot checks) | /assets/charts/ (reproducible scripts)",
             font_size=9, font_color=GRAY)

add_footer(slide, "This slide demonstrates the audit mechanism. Every claim in the deck has a corresponding trace file.")

# ============================================================
# SLIDE 14: Number Reconciliation
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_background(slide, WHITE)

add_text_box(slide, Inches(0.8), Inches(0.3), Inches(11), Inches(0.6),
             "Number Reconciliation: 5 Spot-Checked Numbers",
             font_size=28, font_color=TEAL, bold=True)

add_text_box(slide, Inches(0.8), Inches(0.9), Inches(11), Inches(0.4),
             "Proof that the audit mechanism works. Each number is traced to its source and verified.",
             font_size=14, font_color=GRAY)

# Table
recon_data = [
    ("Revenue FY2026", "$955M", "income_statement_annual.csv", "Total Revenue, 2026-01-31: 955,224,000", "MATCH"),
    ("Op Margin FY2026", "-7.4%", "income_statement_annual.csv", "-70,481,000 / 955,224,000 = -7.38% → -7.4%", "MATCH"),
    ("FCF FY2026", "$222M", "cash_flow_annual.csv", "Free Cash Flow, 2026-01-31: 222,029,000", "MATCH"),
    ("NRR FY2026", "115%", "financial_summary.json", "nrr[4]: 115", "MATCH"),
    ("DCF Base Case", "$38.52", "memo/gitlab_investment_memo.md", "Implied Price: $38.52/share", "MATCH"),
]

# Table header
header_y = Inches(1.4)
col_widths = [Inches(2.0), Inches(1.2), Inches(3.0), Inches(4.5), Inches(1.0)]
col_starts = [Inches(0.5)]
for w in col_widths[:-1]:
    col_starts.append(col_starts[-1] + w)

headers = ["Number", "Deck Value", "Source File", "Source Location", "Verdict"]
for i, (header, start, width) in enumerate(zip(headers, col_starts, col_widths)):
    add_shape_rect(slide, start, header_y, width, Inches(0.4), TEAL)
    add_text_box(slide, start, header_y, width, Inches(0.4),
                 header, font_size=11, font_color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)

# Table rows
for row_idx, (num, deck_val, source, location, verdict) in enumerate(recon_data):
    row_y = Inches(1.8 + row_idx * 0.85)
    bg_color = WHITE if row_idx % 2 == 0 else LIGHT_GRAY
    for col_idx, (val, start, width) in enumerate(zip([num, deck_val, source, location, verdict], col_starts, col_widths)):
        add_shape_rect(slide, start, row_y, width, Inches(0.8), bg_color, RGBColor(0xE5, 0xE7, 0xEB))
        fc = GREEN if col_idx == 4 else DARK_GRAY
        add_text_box(slide, start, row_y, width, Inches(0.8),
                     val, font_size=10, font_color=fc, bold=(col_idx == 4), alignment=PP_ALIGN.CENTER)

add_footer(slide, "Full reconciliation: /audit/reconciliation.md | All numbers: /audit/numbers.md")

# ============================================================
# SLIDE 15: Q&A / Discussion
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_background(slide, WHITE)

add_shape_rect(slide, Inches(0), Inches(0), Inches(13.333), Inches(4), TEAL)

add_text_box(slide, Inches(1), Inches(1.2), Inches(11), Inches(1),
             "Questions & Discussion", font_size=40, font_color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)

add_text_box(slide, Inches(1), Inches(2.2), Inches(11), Inches(0.8),
             "The deck is for discussion, not redelivery. You've read the memo — this is where we talk.",
             font_size=18, font_color=TEAL_LIGHT, alignment=PP_ALIGN.CENTER)

# Discussion prompts
prompts_text = "\n".join([
    "What assumptions in the DCF do you find most questionable?",
    "How confident are you in the NRR trajectory?",
    "What would you do differently if you were the analyst?",
    "Is the audit trail sufficient for your needs?",
    "What additional analysis would you want to see?"
])

add_text_box(slide, Inches(2), Inches(3.5), Inches(9), Inches(3.0),
             prompts_text, font_size=16, font_color=WHITE, alignment=PP_ALIGN.CENTER)

add_text_box(slide, Inches(1), Inches(6.5), Inches(11), Inches(0.5),
             "Recommendation: BUY | Target: $42.00 | Current: $21.51 | Confidence: 7/10",
             font_size=16, font_color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)

add_text_box(slide, Inches(1), Inches(7.0), Inches(11), Inches(0.3),
             "Repo: /workspace/ | Audit: /audit/ | Charts: /assets/charts/ | Traces: /audit/traces/",
             font_size=10, font_color=TEAL_LIGHT, alignment=PP_ALIGN.CENTER)

# Save
os.makedirs('/workspace/deck', exist_ok=True)
prs.save('/workspace/deck/gitlab_board_presentation.pptx')
print("Presentation saved: /workspace/deck/gitlab_board_presentation.pptx")
print(f"Total slides: {len(prs.slides)}")
